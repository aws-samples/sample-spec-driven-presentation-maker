# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""spec-driven-presentation-maker Local MCP Tools (Layer 2).

Security: AWS manages infrastructure security. You manage access control,
data classification, and IAM policies. See SECURITY.md for details.

Each tool wraps the skill/ engine for local filesystem usage.
No AWS dependencies. No deck management (that's Layer 3).
"""

import re
import tempfile
from datetime import datetime
from pathlib import Path
from typing import Any

import json


def _get_output_base_dir() -> Path:
    """Return output base directory from config, with tilde expansion."""
    try:
        from sdpm.config import get_output_dir
        return get_output_dir()
    except Exception:
        return Path.home() / "Documents" / "SDPM-Presentations"


def init_presentation(
    name: str,
    template: str,
    skill_dir: Path,
    style: str = "",
) -> dict[str, Any]:
    """Create a presentation workspace with empty JSON, specs/, and optional template.

    When style is provided, copies the style HTML as specs/art-direction.html.
    When style is empty, creates an empty specs/art-direction.md placeholder.

    Args:
        name: Presentation name (used in directory name).
        skill_dir: Path to the skill/ directory.
        template: Template name (from list-templates) or full path. Required.
        style: Style name (e.g. "elegant-dark"). Empty string means no style.

    Returns:
        Dict with output_dir, json_path, template info, and workspace file list.

    Raises:
        FileNotFoundError: If the specified style does not exist.
    """
    from datetime import datetime as _dt

    ts = _dt.now().strftime("%Y%m%d-%H%M")
    dir_name = f"{ts}-{name}" if name else ts
    out_dir = _get_output_base_dir() / dir_name
    out_dir.mkdir(parents=True, exist_ok=True)

    pres_data: dict[str, Any] = {"fonts": {"fullwidth": None, "halfwidth": None}, "slides": []}

    if template:
        templates_dir = skill_dir / "templates"
        template_src = Path(template).expanduser()
        if not template_src.exists() and "/" not in template and "\\" not in template:
            candidate = templates_dir / (template if template.endswith(".pptx") else template + ".pptx")
            if candidate.exists():
                template_src = candidate
        template_src = template_src.resolve()
        if template_src.exists():
            pres_data["template"] = template_src.name
            from sdpm.analyzer import extract_fonts
            pres_data["fonts"] = extract_fonts(template_src)

    json_path = out_dir / "presentation.json"
    import json
    json_path.write_text(json.dumps(pres_data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")

    # Create specs/ directory with empty spec files
    specs_dir = out_dir / "specs"
    specs_dir.mkdir(exist_ok=True)
    spec_files = ("brief.md", "outline.md")
    if style:
        # Copy style HTML as art-direction.html
        styles_dir = skill_dir / "references" / "examples" / "styles"
        style_name = style if style.endswith(".html") else style + ".html"
        style_src = styles_dir / style_name
        if not style_src.exists():
            raise FileNotFoundError(f"Style not found: {style_src}")
        import shutil
        shutil.copy2(style_src, specs_dir / "art-direction.html")
        art_file = "specs/art-direction.html"
    else:
        (specs_dir / "art-direction.md").touch()
        art_file = "specs/art-direction.md"
    for spec_name in spec_files:
        (specs_dir / spec_name).touch()

    workspace = ["presentation.json"] + [f"specs/{s}" for s in spec_files] + [art_file]

    return {
        "output_dir": str(out_dir),
        "json_path": str(json_path),
        "template": pres_data.get("template", ""),
        "fonts": pres_data.get("fonts", {}),
        "workspace": workspace,
    }


def analyze_template(
    template_path: str,
    skill_dir: Path,
    layout: str = "",
) -> dict[str, Any]:
    """Analyze a PPTX template to extract layouts, colors, fonts.

    Args:
        template_path: Path to .pptx file or template name.
        skill_dir: Path to the skill/ directory.
        layout: Optional layout name or index for detailed placeholder info.

    Returns:
        Dict with layouts, theme_colors, and fonts. If layout specified, includes placeholder details.
    """
    from sdpm.analyzer import analyze_template as _analyze

    if not template_path:
        raise FileNotFoundError("template_path is required. Use start_presentation to see available templates.")

    path = Path(template_path)
    if not path.exists():
        # Try as template name in templates/ directory
        path = skill_dir / "templates" / f"{template_path}.pptx"
    if not path.exists():
        raise FileNotFoundError(f"Template not found: {template_path}")

    output_dir = Path(tempfile.gettempdir()) / "sdpm" / "templates" / path.stem
    result = _analyze(path, output_dir)

    if layout:
        from sdpm.analyzer import get_layout_detail
        detail = get_layout_detail(output_dir, layout)
        if detail:
            result["layout_detail"] = detail
        else:
            result["layout_detail_error"] = f"Layout not found: {layout}"

    return result


def generate_pptx(
    slides_json_path: str,
    skill_dir: Path,
    template: str = "",
    output_path: str = "",
) -> dict[str, Any]:
    """Generate PPTX from a JSON file.

    Args:
        slides_json_path: Path to slides JSON file.
        skill_dir: Path to the skill/ directory.
        template: Template name or path. Empty uses default.
        output_path: Output path. Auto-generated if empty.

    Returns:
        Dict with output_path and slide summary.
    """
    from sdpm.builder import PPTXBuilder, resolve_override, validate_icons_in_json
    from sdpm.utils.io import read_json

    input_path = Path(slides_json_path)
    if not input_path.exists():
        raise FileNotFoundError(f"Slides JSON not found: {slides_json_path}")

    try:
        data = read_json(input_path)
    except (json.JSONDecodeError, ValueError) as e:
        return {"error": f"Invalid JSON in {slides_json_path}: {e}"}
    templates_dir = skill_dir / "templates"

    # Resolve template
    template_file = None
    custom_template = False
    if template:
        p = Path(template)
        if p.exists():
            template_file = p
            custom_template = True
        else:
            named = templates_dir / f"{template}.pptx"
            if named.exists():
                template_file = named
                custom_template = True
            else:
                raise FileNotFoundError(f"Template not found: {template}")
    elif data.get("template"):
        t = input_path.parent / data["template"]
        if t.exists():
            template_file = t
            custom_template = True

    if not template_file:
        template_file = templates_dir / "default.pptx"
        if not template_file.exists():
            raise FileNotFoundError(f"Default template not found: {template_file}")

    # Validate icons
    missing = validate_icons_in_json(data)
    if missing:
        return {"error": f"Missing icons ({len(missing)}): {', '.join(sorted(missing)[:10])}"}

    # Build
    fonts = data.get("fonts")
    builder = PPTXBuilder(
        template_file, custom_template=custom_template, fonts=fonts,
        base_dir=input_path.parent,
        default_text_color=data.get("defaultTextColor"),
    )
    slides = data.get("slides", [])

    id_map = {}
    for s in slides:
        if "id" in s:
            id_map[s["id"]] = s

    for s in slides:
        builder.add_slide(resolve_override(s, id_map))

    # Output path
    if not output_path:
        ts = datetime.now().strftime("%Y%m%d-%H%M%S")
        p = input_path.with_suffix(".pptx")
        out = p.with_stem(f"{p.stem}_{ts}")
    else:
        out = Path(output_path)

    builder.save(out)

    summary = []
    for i, s in enumerate(slides, 1):
        title = s.get("title", "(no title)")
        if isinstance(title, dict):
            title = title.get("text", "(no title)")
        summary.append(f"page{i:02d} - {title}")

    return {
        "output_path": str(out),
        "slide_count": len(slides),
        "slides": summary,
    }


def preview(
    pptx_path: str,
    pages: str = "",
) -> dict[str, Any]:
    """Generate PNG previews of a PPTX file.

    Args:
        pptx_path: Path to .pptx file.
        pages: Comma-separated page numbers. Empty for all.

    Returns:
        Dict with generated PNG paths.
    """
    import platform

    path = Path(pptx_path).resolve()
    if not path.exists():
        raise FileNotFoundError(f"PPTX not found: {pptx_path}")

    os_name = platform.system()
    if os_name == "Darwin":
        return _preview_macos(path, pages)
    elif os_name == "Windows" or _is_wsl():
        return _preview_windows(path, pages)
    else:
        return {
            "error": "Preview requires Microsoft PowerPoint. "
                     "On macOS: install Microsoft PowerPoint. "
                     "On Windows/WSL: PowerPoint must be installed. "
                     "Alternatively, open the generated PPTX file directly."
        }


def _is_wsl() -> bool:
    """Check if running under Windows Subsystem for Linux."""
    proc_version = Path("/proc/version")
    return proc_version.exists() and "microsoft" in proc_version.read_text().lower()


def _preview_macos(pptx_path: Path, pages: str) -> dict[str, Any]:
    """Generate PNG previews on macOS via PowerPoint + AppleScript.

    Args:
        pptx_path: Resolved path to .pptx file.
        pages: Comma-separated page numbers.

    Returns:
        Dict with generated PNG paths.
    """
    import subprocess
    import glob

    from pptx import Presentation
    from sdpm.preview import _mac_open_pptx_background, _mac_restore_pptx_focus_async

    output_dir = Path("/tmp/pptx-preview")
    output_dir.mkdir(parents=True, exist_ok=True)
    pdf_path = output_dir / "slides.pdf"

    # Open in PowerPoint background, export PDF
    pptx_name = pptx_path.name
    restore_info = _mac_open_pptx_background(pptx_path)
    _mac_restore_pptx_focus_async(restore_info)

    import time
    time.sleep(2)

    script = f'''
    tell application "Microsoft PowerPoint"
        set theDoc to presentation "{pptx_name}"
        set outPath to (POSIX file "{pdf_path}") as text
        save theDoc in outPath as save as PDF
        close theDoc saving no
    end tell
    '''
    result = subprocess.run(["osascript", "-e", script], capture_output=True, text=True)  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
    if result.returncode != 0 or not pdf_path.exists():
        return {"error": f"PDF export failed. Is Microsoft PowerPoint installed? Error: {result.stderr}"}

    # PDF → PNG via pdftoppm
    cmd = ["pdftoppm", "-png", str(pdf_path), str(output_dir / "page")]
    result = subprocess.run(cmd, capture_output=True, text=True)  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
    if result.returncode != 0:
        return {"error": f"PNG conversion failed. Is poppler (pdftoppm) installed? Error: {result.stderr}"}

    # Parse pages filter
    pages_set = None
    if pages:
        pages_set = set(int(p.strip()) for p in pages.split(","))

    # Rename with slide titles
    prs = Presentation(str(pptx_path))
    titles = {}
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip().replace("\n", " ")[:30]
        title = re.sub(r'[\\/:*?"<>|]', '', title)
        titles[i] = title or "notitle"

    generated = []
    for png in sorted(glob.glob(str(output_dir / "page-*.png"))):
        match = re.match(r'page-(\d+)\.png', Path(png).name)
        if match:
            num = int(match.group(1))
            if pages_set and num not in pages_set:
                Path(png).unlink()
                continue
            new_name = f"page{num:02d}-{titles.get(num, 'notitle')}.png"
            new_path = output_dir / new_name
            Path(png).rename(new_path)
            generated.append(str(new_path))

    pdf_path.unlink(missing_ok=True)

    return {"preview_dir": str(output_dir), "files": generated}


def _preview_windows(pptx_path: Path, pages: str) -> dict[str, Any]:
    """Generate PNG previews on Windows/WSL via PowerShell COM.

    Args:
        pptx_path: Resolved path to .pptx file.
        pages: Comma-separated page numbers.

    Returns:
        Dict with generated PNG paths.
    """
    import subprocess
    import glob

    from pptx import Presentation

    output_dir = pptx_path.parent / "preview"
    output_dir.mkdir(parents=True, exist_ok=True)
    png_dir = output_dir / "ppt_png"
    png_dir.mkdir(parents=True, exist_ok=True)

    if _is_wsl():
        win_pptx = subprocess.run(["wslpath", "-w", str(pptx_path)], capture_output=True, text=True).stdout.strip()  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
        win_png = subprocess.run(["wslpath", "-w", str(png_dir)], capture_output=True, text=True).stdout.strip()  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
        ps_cmd = (
            f"$app = New-Object -ComObject PowerPoint.Application; "
            f"$prs = $app.Presentations.Open('{win_pptx}', "
            f"[Microsoft.Office.Core.MsoTriState]::msoTrue, "
            f"[Microsoft.Office.Core.MsoTriState]::msoFalse, "
            f"[Microsoft.Office.Core.MsoTriState]::msoFalse); "
            f"$prs.SaveAs('{win_png}', 18); "
            f"$prs.Close(); $app.Quit()"
        )
        shell = "/mnt/c/Windows/System32/WindowsPowerShell/v1.0/powershell.exe"
    else:
        ps_cmd = (
            f"$app = New-Object -ComObject PowerPoint.Application; "
            f"$prs = $app.Presentations.Open('{pptx_path}', "
            f"[Microsoft.Office.Core.MsoTriState]::msoTrue, "
            f"[Microsoft.Office.Core.MsoTriState]::msoFalse, "
            f"[Microsoft.Office.Core.MsoTriState]::msoFalse); "
            f"$prs.SaveAs('{png_dir}', 18); "
            f"$prs.Close(); $app.Quit()"
        )
        shell = "powershell.exe"

    result = subprocess.run([shell, "-Command", ps_cmd], capture_output=True, timeout=60)  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
    if result.returncode != 0:
        stderr = result.stderr.decode("cp932", errors="replace") if result.stderr else ""
        return {"error": f"PNG export failed: {stderr}"}

    pages_set = None
    if pages:
        pages_set = set(int(p.strip()) for p in pages.split(","))

    prs = Presentation(str(pptx_path))
    titles = {}
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip().replace("\n", " ")[:30]
        title = re.sub(r'[\\/:*?"<>|]', '', title)
        titles[i] = title or "notitle"

    png_files = sorted(glob.glob(str(png_dir / "**" / "*.png"), recursive=True))
    if not png_files:
        png_files = sorted(glob.glob(str(png_dir / "**" / "*.PNG"), recursive=True))

    generated = []
    for idx, png in enumerate(png_files, 1):
        if pages_set and idx not in pages_set:
            Path(png).unlink()
            continue
        new_name = f"page{idx:02d}-{titles.get(idx, 'notitle')}.png"
        new_path = output_dir / new_name
        Path(png).rename(new_path)
        generated.append(str(new_path))

    import shutil
    shutil.rmtree(png_dir, ignore_errors=True)

    return {"preview_dir": str(output_dir), "files": generated}


def search_assets(
    query: str,
    skill_dir: Path,
    limit: int = 20,
    source_filter: str = "",
    type_filter: str = "",
    theme_filter: str = "",
) -> dict[str, Any]:
    """Search assets (icons, images, etc.) by keyword using local manifests.

    Args:
        query: Search keyword.
        skill_dir: Path to the skill/ directory.
        limit: Maximum results.
        source_filter: Filter by source name (e.g. "aws", "material").
        type_filter: Filter by asset type.
        theme_filter: Filter by theme.

    Returns:
        Dict with matching assets.
    """
    from sdpm.assets import search_assets as _search

    return {
        "query": query,
        "results": _search(
            query, limit=limit, source_filter=source_filter or None,
            type_filter=type_filter or None, theme_filter=theme_filter or None,
        ),
    }


def list_asset_sources(skill_dir: Path) -> dict[str, Any]:
    """List available asset sources with counts and descriptions.

    Args:
        skill_dir: Path to the skill/ directory.

    Returns:
        Dict with list of sources.
    """
    from sdpm.assets import list_sources

    return {"sources": list_sources()}


def _list_category(category_dir: Path) -> list[dict[str, str]]:
    """List all .md, .pptx, and .html files in a category directory with descriptions.

    For .md files, extracts description from frontmatter.
    For .pptx files, extracts description from first slide's speaker notes.
    For .html files, extracts description from <title> tag.
    Recurses into subdirectories (e.g. styles/).

    Args:
        category_dir: Path to the category directory.

    Returns:
        List of dicts with name, description, and path.
    """
    from sdpm.reference import _get_description

    if not category_dir.exists():
        return []
    items = []
    seen: set[str] = set()
    for f in sorted(category_dir.rglob("*")):
        if f.suffix not in (".md", ".pptx", ".html"):
            continue
        # Use relative path from category_dir as the unique key
        rel = f.relative_to(category_dir)
        stem_key = str(rel.with_suffix(""))
        if stem_key in seen:
            continue
        seen.add(stem_key)
        desc = _get_description(f)
        items.append({
            "name": stem_key,
            "description": desc,
            "path": f"{category_dir.name}/{rel}",
        })
    return items


def _read_docs(category_dir: Path, names: list[str]) -> list[dict[str, Any]]:
    """Read one or more documents from a category directory.

    Supports .md (returned as text), .pptx (rendered via get_pptx_notes),
    and .html (returned as text).
    Names can include page specifiers for pptx: "common/patterns/3" or "common/patterns/all".

    Args:
        category_dir: Path to the category directory.
        names: List of document names. Supports "name", "name/page", "name/all".

    Returns:
        List of dicts with name, path, and content.
    """
    results = []
    for name in names:
        # Parse page specifier from name (e.g. "common/patterns/3" → file="common/patterns", pages=[3])
        pages = None
        has_page_specifier = False
        parts = name.rsplit("/", 1)
        file_name = name
        if len(parts) == 2 and (parts[1].isdigit() or parts[1] == "all"):
            file_name = parts[0]
            has_page_specifier = True
            pages = None if parts[1] == "all" else [int(parts[1])]

        # Try .md first, then .pptx, then .html
        md_path = category_dir / f"{file_name}.md"
        pptx_path = category_dir / f"{file_name}.pptx"
        html_path = category_dir / f"{file_name}.html"

        if md_path.exists():
            results.append({
                "name": file_name,
                "path": f"{category_dir.name}/{file_name}.md",
                "content": md_path.read_text(encoding="utf-8"),
            })
        elif pptx_path.exists():
            if not has_page_specifier:
                # No page specifier — return slide description listing
                from sdpm.reference import list_pptx_descriptions
                descriptions = list_pptx_descriptions(str(pptx_path))
                content = "\n".join(f"  {page:>3}  {desc}" for page, desc in descriptions)
                results.append({
                    "name": name,
                    "path": f"{category_dir.name}/{file_name}.pptx",
                    "content": content,
                })
            else:
                from sdpm.reference import get_pptx_notes
                notes = get_pptx_notes(pptx_path, pages=pages)
                content_lines: list[str] = []
                for page_num, note_text in notes:
                    content_lines.append(f"## Page {page_num}\n\n{note_text}\n")
                results.append({
                    "name": name,
                    "path": f"{category_dir.name}/{file_name}.pptx",
                    "content": "\n".join(content_lines),
                })
        elif html_path.exists():
            results.append({
                "name": file_name,
                "path": f"{category_dir.name}/{file_name}.html",
                "content": html_path.read_text(encoding="utf-8"),
            })
        else:
            available = sorted({
                str(f.relative_to(category_dir).with_suffix(""))
                for f in category_dir.rglob("*")
                if f.suffix in (".md", ".pptx", ".html")
            })
            raise FileNotFoundError(
                f"'{file_name}' not found in {category_dir.name}/. Available: {', '.join(available)}"
            )
    return results


def list_styles(skill_dir: Path) -> dict[str, Any]:
    """List available design styles from references/examples/styles/.

    Extracts name and description from each HTML file's <title> tag.

    Args:
        skill_dir: Path to the skill/ directory.

    Returns:
        Dict with styles list (name + description).
    """
    styles_dir = skill_dir / "references" / "examples" / "styles"
    if not styles_dir.exists():
        return {"styles": []}

    styles: list[dict[str, str]] = []
    for f in sorted(styles_dir.iterdir()):
        if f.suffix != ".html" or f.name.startswith("."):
            continue
        description = ""
        try:
            content = f.read_text(encoding="utf-8")
            m = re.search(r"<title>(.*?)</title>", content, re.IGNORECASE)
            if m:
                description = m.group(1).strip()
        except Exception:
            pass
        styles.append({"name": f.stem, "description": description})
    return {"styles": styles}


def read_examples(names: list[str], skill_dir: Path) -> dict[str, Any]:
    """Read design examples (components/patterns).

    Without page specifier returns a listing of slide descriptions.
    With page specifier returns full content.

    Args:
        names: Example names (e.g. ["patterns", "patterns/3", "components/all"]).
        skill_dir: Path to the skill/ directory.

    Returns:
        Dict with documents list.
    """
    return {"documents": _read_docs(skill_dir / "references" / "examples", names)}


def list_workflows(skill_dir: Path) -> dict[str, Any]:
    """List all workflow documents.

    Args:
        skill_dir: Path to the skill/ directory.

    Returns:
        Dict with items list.
    """
    return {"items": _list_category(skill_dir / "references" / "workflows")}


def read_workflows(names: list[str], skill_dir: Path) -> dict[str, Any]:
    """Read one or more workflow documents.

    Args:
        names: Workflow names (e.g. ["create-new-2-build", "slide-json-spec"]).
        skill_dir: Path to the skill/ directory.

    Returns:
        Dict with documents list.
    """
    return {"documents": _read_docs(skill_dir / "references" / "workflows", names)}


def list_guides(skill_dir: Path) -> dict[str, Any]:
    """List all guide documents.

    Args:
        skill_dir: Path to the skill/ directory.

    Returns:
        Dict with items list.
    """
    return {"items": _list_category(skill_dir / "references" / "guides")}


def read_guides(names: list[str], skill_dir: Path) -> dict[str, Any]:
    """Read one or more guide documents.

    Args:
        names: Guide names (e.g. ["design-rules"]).
        skill_dir: Path to the skill/ directory.

    Returns:
        Dict with documents list.
    """
    return {"documents": _read_docs(skill_dir / "references" / "guides", names)}


def code_block(
    code: str,
    language: str = "python",
    theme: str = "dark",
    x: int = 0,
    y: int = 0,
    width: int = 800,
    height: int = 300,
) -> list[dict[str, Any]]:
    """Generate slide elements JSON for a syntax-highlighted code block.

    Args:
        code: Source code text.
        language: Programming language.
        theme: Color theme ("dark" or "light").
        x: X position in pixels.
        y: Y position in pixels.
        width: Width in pixels.
        height: Height in pixels.

    Returns:
        List of slide element dicts for the code block.
    """
    from sdpm.utils.text import highlight_code
    from sdpm.builder.constants import CODE_COLORS

    colors = CODE_COLORS.get(theme, CODE_COLORS["dark"])
    label_height = 22

    # Background
    elements = [
        {
            "type": "shape",
            "shape_type": "rectangle",
            "x": x, "y": y, "width": width, "height": height,
            "fill": colors["bg"],
            "corner_radius": 8,
        },
        {
            "type": "shape",
            "shape_type": "rectangle",
            "x": x, "y": y, "width": width, "height": label_height,
            "fill": colors["label_bg"],
            "corner_radius": 8,
        },
        {
            "type": "textbox",
            "x": x + 8, "y": y + 2,
            "width": width - 16, "height": label_height - 4,
            "text": [{"text": language, "font_size": 9, "color": colors["label_fg"]}],
        },
    ]

    # Highlighted code
    spans = highlight_code(code, language, theme)
    elements.append({
        "type": "textbox",
        "x": x + 12, "y": y + label_height + 4,
        "width": width - 24, "height": height - label_height - 8,
        "text": spans,
    })

    return elements


def pptx_to_json(
    pptx_path: str,
) -> dict[str, Any]:
    """Convert PPTX to JSON representation.

    Args:
        pptx_path: Path to .pptx file.

    Returns:
        Dict with JSON representation of slides.
    """
    from sdpm.converter import pptx_to_json as _convert

    path = Path(pptx_path)
    if not path.exists():
        raise FileNotFoundError(f"PPTX not found: {pptx_path}")

    return _convert(str(path))
