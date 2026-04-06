# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Presentation app backend abstraction — PowerPoint / LibreOffice."""

import os
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path




class PresentationBackend:
    """Abstract base for presentation app operations."""

    name: str = "none"

    def export_pdf(self, pptx_path: Path, pdf_path: Path) -> bool:
        """Export PPTX to PDF. Returns True on success."""
        raise NotImplementedError

    def refresh_autofit(self, pptx_path: Path, pdf_path: Path | None = None) -> bool:
        """Refresh autofit and optionally export PDF. Returns True on success."""
        raise NotImplementedError

    def refresh_autofit_for_convert(self, pptx_path: Path) -> Path | None:
        """Re-save PPTX with autofit applied for converter. Returns temp path or None."""
        raise NotImplementedError

    def open_background(self, pptx_path: Path):
        """Open PPTX in background for interactive use. Returns restore info."""
        return None


# ---------------------------------------------------------------------------
# PowerPoint
# ---------------------------------------------------------------------------

def _is_wsl() -> bool:
    return Path("/proc/version").exists() and "microsoft" in Path("/proc/version").read_text().lower()


class PowerPointBackend(PresentationBackend):
    """PowerPoint via AppleScript (macOS) or COM (Windows/WSL)."""

    name = "powerpoint"

    def export_pdf(self, pptx_path: Path, pdf_path: Path) -> bool:
        if sys.platform == "darwin":
            return self._mac_export_pdf(pptx_path, pdf_path)
        elif sys.platform == "win32" or _is_wsl():
            return self._win_export_pdf(pptx_path, pdf_path)
        return False

    def refresh_autofit(self, pptx_path: Path, pdf_path: Path | None = None) -> bool:
        if sys.platform == "darwin":
            return self._mac_refresh_autofit(pptx_path, pdf_path)
        elif sys.platform == "win32":
            return self._win_refresh_autofit(pptx_path, pdf_path)
        elif _is_wsl():
            return self._wsl_refresh_autofit(pptx_path, pdf_path)
        return False

    def refresh_autofit_for_convert(self, pptx_path: Path) -> Path | None:
        from pptx import Presentation
        from pptx.oxml.ns import qn
        prs = Presentation(str(pptx_path))
        targets = []
        for si, slide in enumerate(prs.slides, 1):
            for shi, shape in enumerate(slide.shapes, 1):
                sp = shape._element
                nvSpPr = sp.find(qn('p:nvSpPr'))
                if nvSpPr is None:
                    continue
                cNvSpPr = nvSpPr.find(qn('p:cNvSpPr'))
                if cNvSpPr is None or cNvSpPr.get('txBox') != '1':
                    continue
                targets.append((si, shi))
        if not targets:
            return None
        tmp_dir = tempfile.mkdtemp()
        tmp_pptx = Path(tmp_dir) / pptx_path.name
        shutil.copy2(pptx_path, tmp_pptx)
        if sys.platform == "darwin":
            cmd = self._mac_convert_autofit_cmd(tmp_pptx, targets)
        elif sys.platform == "win32" or _is_wsl():
            cmd = self._win_convert_autofit_cmd(tmp_pptx, targets)
        else:
            shutil.rmtree(tmp_dir)
            return None
        try:
            result = subprocess.run(cmd, capture_output=True, timeout=60)  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
            if result.returncode == 0:
                return tmp_pptx
        except (subprocess.TimeoutExpired, FileNotFoundError):
            pass
        shutil.rmtree(tmp_dir)
        return None

    def open_background(self, pptx_path: Path):
        if sys.platform == "darwin":
            return _mac_open_pptx_background(pptx_path)
        return None

    # --- macOS AppleScript helpers ---

    def _mac_refresh_autofit(self, pptx_path: Path, pdf_path: Path | None) -> bool:
        pdf_lines = ""
        if pdf_path:
            pdf_lines = f'''
                set outPath to (POSIX file "{pdf_path}") as text
                save pres in outPath as save as PDF
            '''
        pptx_name = pptx_path.name
        restore_info = _mac_open_pptx_background(pptx_path)
        _mac_restore_pptx_focus_async(restore_info)
        import time
        time.sleep(2)
        script = f'''
            tell application "Microsoft PowerPoint"
                set pres to presentation "{pptx_name}"
                set slideCount to count of slides of pres
                if (count of shapes of slide 1 of pres) > 0 then
                    set sh to shape 1 of slide 1 of pres
                    set w to width of sh
                    set width of sh to w + 1
                    set width of sh to w
                end if
                set waitTime to 1 + slideCount * 0.3
                delay waitTime
                save pres
                {pdf_lines}
                close pres
            end tell
        '''
        return _run_cmd(["osascript", "-e", script])

    def _mac_export_pdf(self, pptx_path: Path, pdf_path: Path) -> bool:
        pptx_name = pptx_path.name
        restore_info = _mac_open_pptx_background(pptx_path)
        _mac_restore_pptx_focus_async(restore_info)
        import time
        time.sleep(2)
        script = f'''
            tell application "Microsoft PowerPoint"
                set theDoc to presentation "{pptx_name}"
                set outPath to (POSIX file "{pdf_path}") as text
                save theDoc in outPath as save as PDF
                close theDoc saving no
            end tell
        '''
        return _run_cmd(["osascript", "-e", script])

    def _mac_convert_autofit_cmd(self, tmp_pptx: Path, targets: list) -> list:
        resize_lines = "\n".join(
            f'                try\n'
            f'                    tell text frame of shape {shi} of slide {si} of pres\n'
            f'                        set auto size to shape to fit text\n'
            f'                    end tell\n'
            f'                end try'
            for si, shi in targets
        )
        pptx_name = tmp_pptx.name
        restore_info = _mac_open_pptx_background(tmp_pptx)
        _mac_restore_pptx_focus_async(restore_info)
        import time
        time.sleep(2)
        script = f'''
    tell application "Microsoft PowerPoint"
        set pres to presentation "{pptx_name}"
{resize_lines}
        delay 1
        save pres
        close pres
    end tell
'''
        return ["osascript", "-e", script]

    # --- Windows/WSL helpers ---

    def _win_refresh_autofit(self, pptx_path: Path, pdf_path: Path | None) -> bool:
        pdf_line = ""
        if pdf_path:
            pdf_line = f"$prs.SaveAs('{pdf_path}', 32); "
        ps_cmd = (
            f"$app = New-Object -ComObject PowerPoint.Application; "
            f"$prs = $app.Presentations.Open('{pptx_path}'); "
            f"$s = $prs.Slides[1]; "
            f"if ($s.Shapes.Count -gt 0) {{ $sh = $s.Shapes[1]; $w = $sh.Width; $sh.Width = $w + 1; $sh.Width = $w }}; "
            f"$prs.Save(); {pdf_line}$prs.Close(); $app.Quit()"
        )
        return _run_cmd(["powershell.exe", "-Command", ps_cmd])

    def _wsl_refresh_autofit(self, pptx_path: Path, pdf_path: Path | None) -> bool:
        win_path = _wsl_path(pptx_path)
        pdf_line = ""
        if pdf_path:
            win_pdf = _wsl_path(pdf_path)
            pdf_line = f"$prs.SaveAs('{win_pdf}', 32); "
        ps_cmd = (
            f"$app = New-Object -ComObject PowerPoint.Application; "
            f"$prs = $app.Presentations.Open('{win_path}'); "
            f"$s = $prs.Slides[1]; "
            f"if ($s.Shapes.Count -gt 0) {{ $sh = $s.Shapes[1]; $w = $sh.Width; $sh.Width = $w + 1; $sh.Width = $w }}; "
            f"$prs.Save(); {pdf_line}$prs.Close(); $app.Quit()"
        )
        return _run_cmd(["/mnt/c/Windows/System32/WindowsPowerShell/v1.0/powershell.exe", "-Command", ps_cmd])

    def _win_export_pdf(self, pptx_path: Path, pdf_path: Path) -> bool:
        if _is_wsl():
            win_pptx = _wsl_path(pptx_path)
            win_pdf = _wsl_path(pdf_path)
        else:
            win_pptx, win_pdf = str(pptx_path), str(pdf_path)
        ps_cmd = (
            f"$app = New-Object -ComObject PowerPoint.Application; "
            f"$prs = $app.Presentations.Open('{win_pptx}', "
            f"[Microsoft.Office.Core.MsoTriState]::msoTrue, "
            f"[Microsoft.Office.Core.MsoTriState]::msoFalse, "
            f"[Microsoft.Office.Core.MsoTriState]::msoFalse); "
            f"$prs.SaveAs('{win_pdf}', 32); "
            f"$prs.Close(); $app.Quit()"
        )
        shell = ["/mnt/c/Windows/System32/WindowsPowerShell/v1.0/powershell.exe"] if _is_wsl() else ["powershell.exe"]
        return _run_cmd(shell + ["-Command", ps_cmd])

    def _win_convert_autofit_cmd(self, tmp_pptx: Path, targets: list) -> list:
        if _is_wsl():
            win_path = _wsl_path(tmp_pptx)
        else:
            win_path = str(tmp_pptx)
        resize_lines = "; ".join(
            f"$prs.Slides[{si}].Shapes[{shi}].TextFrame.AutoSize = 1"
            for si, shi in targets
        )
        ps_cmd = (
            f"$app = New-Object -ComObject PowerPoint.Application; "
            f"$prs = $app.Presentations.Open('{win_path}'); "
            f"{resize_lines}; "
            f"$prs.Save(); $prs.Close(); $app.Quit()"
        )
        if _is_wsl():
            return ["/mnt/c/Windows/System32/WindowsPowerShell/v1.0/powershell.exe", "-Command", ps_cmd]
        return ["powershell.exe", "-Command", ps_cmd]


# ---------------------------------------------------------------------------
# LibreOffice
# ---------------------------------------------------------------------------

class LibreOfficeBackend(PresentationBackend):
    """LibreOffice headless."""

    name = "libreoffice"

    def export_pdf(self, pptx_path: Path, pdf_path: Path) -> bool:
        tmp_dir = tempfile.mkdtemp()
        try:
            env = os.environ.copy()
            env["HOME"] = tmp_dir
            subprocess.run(  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
                ["soffice", "--headless", "--convert-to", "pdf", "--outdir", tmp_dir, str(pptx_path)],
                env=env, capture_output=True, text=True, timeout=120, check=True,
            )
            tmp_pdf = Path(tmp_dir) / (pptx_path.stem + ".pdf")
            if tmp_pdf.exists():
                shutil.move(str(tmp_pdf), str(pdf_path))
                return True
            return False
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired, FileNotFoundError):
            return False
        finally:
            shutil.rmtree(tmp_dir, ignore_errors=True)

    def refresh_autofit(self, pptx_path: Path, pdf_path: Path | None = None) -> bool:
        from .autofit import extract_scaling, unlock_autofit
        tmp_dir = tempfile.mkdtemp()
        try:
            env = os.environ.copy()
            env["HOME"] = tmp_dir
            lo_out = Path(tmp_dir) / "resave"
            lo_out.mkdir()
            subprocess.run(  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
                ["soffice", "--headless", "--convert-to", "pptx", "--outdir", str(lo_out), str(pptx_path)],
                env=env, capture_output=True, text=True, timeout=120, check=True,
            )
            resaved = lo_out / pptx_path.name
            if resaved.exists():
                scaling = extract_scaling(resaved)
                unlock_autofit(pptx_path, scaling)
            if pdf_path:
                return self.export_pdf(pptx_path, pdf_path)
            return True
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired, FileNotFoundError):
            return False
        finally:
            shutil.rmtree(tmp_dir, ignore_errors=True)

    def refresh_autofit_for_convert(self, pptx_path: Path) -> Path | None:
        """Re-save via LibreOffice to compute autofit values."""
        tmp_dir = tempfile.mkdtemp()
        tmp_pptx = Path(tmp_dir) / pptx_path.name
        shutil.copy2(pptx_path, tmp_pptx)
        try:
            env = os.environ.copy()
            env["HOME"] = tmp_dir
            lo_out = Path(tmp_dir) / "resave"
            lo_out.mkdir()
            subprocess.run(  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
                ["soffice", "--headless", "--convert-to", "pptx", "--outdir", str(lo_out), str(tmp_pptx)],
                env=env, capture_output=True, text=True, timeout=120, check=True,
            )
            resaved = lo_out / tmp_pptx.name
            if resaved.exists():
                return resaved
        except (subprocess.CalledProcessError, subprocess.TimeoutExpired, FileNotFoundError):
            pass
        shutil.rmtree(tmp_dir, ignore_errors=True)
        return None


# ---------------------------------------------------------------------------
# Detection
# ---------------------------------------------------------------------------

def _has_powerpoint() -> bool:
    if sys.platform == "darwin":
        return Path("/Applications/Microsoft PowerPoint.app").exists()
    elif sys.platform == "win32" or _is_wsl():
        return True  # Assume available; will fail gracefully if not
    return False


def _has_libreoffice() -> bool:
    return shutil.which("soffice") is not None


def detect_backend() -> PresentationBackend | None:
    """Detect available presentation backend.

    Priority: config file > PowerPoint > LibreOffice > None.
    """
    from sdpm.config import get_preview_config

    forced = get_preview_config().get("backend", "").lower()
    if forced == "powerpoint":
        return PowerPointBackend()
    if forced == "libreoffice":
        return LibreOfficeBackend()
    if _has_powerpoint():
        return PowerPointBackend()
    if _has_libreoffice():
        return LibreOfficeBackend()
    return None


# ---------------------------------------------------------------------------
# Shared helpers (moved from preview/__init__.py)
# ---------------------------------------------------------------------------

def _mac_open_pptx_background(pptx_path: Path) -> str:
    """Open PPTX in PowerPoint without stealing fullscreen space."""
    script = '''
    tell application "System Events"
        tell application process "Microsoft PowerPoint"
            set fsName to ""
            set topName to ""
            if (count of windows) > 0 then
                set topName to name of window 1
            end if
            repeat with w in windows
                if value of attribute "AXFullScreen" of w is true then
                    set fsName to name of w
                    set value of attribute "AXFullScreen" of w to false
                    exit repeat
                end if
            end repeat
        end tell
    end tell
    do shell script "open -gF -a 'Microsoft PowerPoint' " & quoted form of "%s"
    return fsName & "|" & topName
    ''' % str(pptx_path)
    result = subprocess.run(["osascript", "-e", script], capture_output=True, text=True)  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
    return result.stdout.strip() if result.returncode == 0 else "|"


def _mac_restore_pptx_focus_async(restore_info: str) -> None:
    """Restore PowerPoint window focus asynchronously."""
    parts = restore_info.split("|", 1)
    fs_name = parts[0] if len(parts) > 0 else ""
    top_name = parts[1] if len(parts) > 1 else ""
    script = _mac_restore_applescript(fs_name, top_name)
    if script:
        subprocess.Popen(["osascript", "-e", script], stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit


def _mac_restore_applescript(fs_name: str, top_name: str) -> str:
    if fs_name:
        return f'''
        tell application "System Events"
            tell application process "Microsoft PowerPoint"
                repeat with w in windows
                    if name of w is "{fs_name}" then
                        set value of attribute "AXFullScreen" of w to true
                        exit repeat
                    end if
                end repeat
            end tell
        end tell
        tell application "Microsoft PowerPoint" to activate
        '''
    elif top_name:
        return f'''
        tell application "System Events"
            tell application process "Microsoft PowerPoint"
                repeat with w in windows
                    if name of w is "{top_name}" then
                        perform action "AXRaise" of w
                        exit repeat
                    end if
                end repeat
            end tell
        end tell
        '''
    return ""


def _run_cmd(cmd: list, timeout: int = 60) -> bool:
    """Run a command, return True on success."""
    try:
        result = subprocess.run(cmd, capture_output=True, timeout=timeout)  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
        if result.returncode == 0:
            return True
        stderr = result.stderr
        if isinstance(stderr, bytes):
            stderr = stderr.decode("utf-8", errors="replace")
        print(f"Warning: Command failed: {stderr.strip()}", file=sys.stderr)
        return False
    except (subprocess.TimeoutExpired, FileNotFoundError, UnicodeDecodeError):
        return False


def _wsl_path(path: Path) -> str:
    """Convert POSIX path to Windows path via wslpath."""
    return subprocess.run(  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
        ["wslpath", "-w", str(path)], capture_output=True, text=True
    ).stdout.strip()
