# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Autofit handling — reads LibreOffice-computed scaling and replaces normAutofit.

After LibreOffice opens and saves a PPTX, normAutofit contains fontScale and
lnSpcReduction values. This module reads those values from the LibreOffice
re-saved copy for reporting, then replaces normAutofit with noAutofit in the
original PPTX without modifying font sizes or spacing.

Ported from hobby pptx-maker png-worker/handlers/autofit.py.
"""

import logging
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn

logger = logging.getLogger("sdpm.preview.autofit")

_NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


def extract_scaling(pptx_path: Path) -> dict[tuple[int, int], tuple[float, float]]:
    """Extract fontScale and lnSpcReduction from normAutofit elements.

    Args:
        pptx_path: Path to the LibreOffice re-saved PPTX.

    Returns:
        Dict mapping (slide_index, shape_index) to (fontScale, lnSpcReduction).
        Only includes shapes where scaling was applied.
    """
    prs = Presentation(str(pptx_path))
    scaling: dict[tuple[int, int], tuple[float, float]] = {}
    for si, slide in enumerate(prs.slides):
        for shi, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            body_pr = shape.text_frame._txBody.find(".//a:bodyPr", _NS)
            if body_pr is None:
                continue
            norm = body_pr.find("a:normAutofit", _NS)
            if norm is None:
                continue
            font_scale = int(norm.get("fontScale", "100000")) / 100000.0
            ln_spc_reduction = int(norm.get("lnSpcReduction", "0")) / 100000.0
            if font_scale < 1.0 or ln_spc_reduction > 0:
                scaling[(si, shi)] = (font_scale, ln_spc_reduction)
    return scaling


def unlock_autofit(original_path: Path, scaling: dict[tuple[int, int], tuple[float, float]]) -> bool:
    """Replace normAutofit with noAutofit without modifying font sizes or spacing.

    Args:
        original_path: Path to the original PPTX file (modified in-place).
        scaling: Dict from extract_scaling() for logging.

    Returns:
        True if any shapes were modified.
    """
    prs = Presentation(str(original_path))
    changed = False
    for si, slide in enumerate(prs.slides):
        for shi, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            body_pr = shape.text_frame._txBody.find(".//a:bodyPr", _NS)
            if body_pr is None:
                continue
            norm = body_pr.find("a:normAutofit", _NS)
            if norm is None:
                continue
            body_pr.remove(norm)
            etree.SubElement(body_pr, qn("a:noAutofit"))
            changed = True
            key = (si, shi)
            if key in scaling:
                font_scale, ln_spc_reduction = scaling[key]
                logger.info(
                    "Unlocked autofit: slide=%d shape=%d fontScale=%.0f%% lnSpcReduction=%.0f%%",
                    si, shi, font_scale * 100, ln_spc_reduction * 100,
                )
    if changed:
        prs.save(str(original_path))
    return changed
