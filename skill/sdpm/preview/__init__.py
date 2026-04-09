# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Preview: autofit refresh, PNG generation, layout imbalance check."""

# Security: AWS manages infrastructure security. You manage access control,
# data classification, and IAM policies. See SECURITY.md for details.
import sys
import tempfile
from pathlib import Path

from pptx import Presentation

from .backend import _is_wsl as _is_wsl, detect_backend


def refresh_autofit(pptx_path, pdf_path=None, pre_autofit_pdf=None):
    """Refresh autofit via detected presentation backend."""
    backend = detect_backend()
    if backend is None:
        print("Warning: Autofit refresh skipped (no presentation app available)", file=sys.stderr)
        return False
    result = backend.refresh_autofit(Path(pptx_path), Path(pdf_path) if pdf_path else None, pre_autofit_pdf=Path(pre_autofit_pdf) if pre_autofit_pdf else None)
    if not result:
        print(f"Warning: Autofit refresh failed ({backend.name})", file=sys.stderr)
    return result


def export_pdf(pptx_path, pdf_path):
    """Export PPTX to PDF via detected presentation backend."""
    backend = detect_backend()
    if backend is None:
        print("Warning: PDF export skipped (no presentation app available)", file=sys.stderr)
        return False
    return backend.export_pdf(Path(pptx_path), Path(pdf_path))


def refresh_autofit_for_convert(pptx_path):
    """Re-save PPTX with autofit applied for converter. Returns temp path or None."""
    backend = detect_backend()
    if backend is None:
        return None
    return backend.refresh_autofit_for_convert(Path(pptx_path))


def check_layout_imbalance(pptx_path, slide_defs=None):
    """Detect slides where bbox centroid deviates from content area center and print results."""
    alerts = check_layout_imbalance_data(pptx_path, slide_defs)
    if alerts:
        print(f"⚠️  Layout bias detected ({len(alerts)} slides):")
        for a in alerts:
            print(f"  page{a['slide']:02d} ({a['layout']}) | {a['bbox']} | centroid offset: {a['offset']} ({a['direction']})")
        print("  → MUST FIX unless the layout type is intentionally asymmetric (e.g. title, section, agenda, thankyou).")
        print("  → Action: Increase element heights, expand spacing between elements, or add content to fill the empty area. Fix one slide at a time. Do NOT batch-fix.")


def check_layout_imbalance_data(pptx_path, slide_defs=None):
    """Detect slides where bbox centroid deviates from content area center."""
    _THRESHOLD = 0.03
    prs = Presentation(str(pptx_path))
    _SW = int(prs.slide_width / 6350)
    _SH = int(prs.slide_height / 6350)
    emu = 6350
    _TITLE_BOTTOM = int(_SH * 0.13)
    _CONTENT_BOTTOM = int(_SH * 0.88)
    _CY = (_TITLE_BOTTOM + _CONTENT_BOTTOM) / 2
    _CA_H = _CONTENT_BOTTOM - _TITLE_BOTTOM
    alerts = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        layout = "content"
        if slide_defs and slide_idx <= len(slide_defs):
            layout = slide_defs[slide_idx - 1].get("layout", "content")
        min_x, min_y, max_x, max_y = _SW, _SH, 0, 0
        has_elem = False
        for shape in slide.shapes:
            if shape.is_placeholder:
                continue
            x = int(shape.left / emu)
            y = int(shape.top / emu)
            min_x, min_y = min(min_x, x), min(min_y, y)
            max_x, max_y = max(max_x, x + int(shape.width / emu)), max(max_y, y + int(shape.height / emu))
            has_elem = True
        if not has_elem:
            continue
        cy = (min_y + max_y) / 2
        dy = (cy - _CY) / _CA_H
        if abs(dy) > _THRESHOLD:
            alerts.append({
                "slide": slide_idx,
                "layout": layout,
                "bbox": f"x={min_x}..{max_x} y={min_y}..{max_y} (of {_SW}x{_SH})",
                "offset": f"{dy:+.1%}",
                "direction": "bottom-heavy" if dy > 0 else "top-heavy",
            })
    return alerts


def unlock_height_constraints(pptx_path):
    """Replace normAutofit with noAutofit and warn on text overflow."""
    from lxml import etree
    prs = Presentation(str(pptx_path))
    nsmap = {'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'}
    changed = False
    warnings = []
    for slide_idx, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            body_pr = shape.text_frame._txBody.find('.//a:bodyPr', nsmap)
            if body_pr is None:
                continue
            norm = body_pr.find('a:normAutofit', nsmap)
            if norm is None:
                continue
            font_scale = int(norm.get('fontScale', '100000')) / 100000.0
            if font_scale < 1.0:
                text_preview = shape.text_frame.text[:40].replace('\n', ' ')
                warnings.append(f"  page{slide_idx:02d} | \"{text_preview}\" | would need fontSize {font_scale:.0%} to fit")
            body_pr.remove(norm)
            etree.SubElement(body_pr, f'{{{nsmap["a"]}}}noAutofit')
            changed = True
    if warnings:
        print(f"⚠️  Text may not fit ({len(warnings)} shapes) — would need fontSize reduction to fit within shape bounds:")
        for line in warnings:
            print(line)
        print("  fontSize reduction also changes text wrapping — exact cause of overflow varies.")
        print("  Check preview for: text overflowing shape bounds, unintended wrapping, overlap with nearby elements.")
        print("  If no visible problem, no action needed.")
        print("  If text overflows, adjust layout (width/height) or content (text length).")
        print("  fontSize has recommended minimums — prefer layout/content adjustments over reducing fontSize.")
        print("  Review surrounding layout together — don't fix in isolation.")
    if changed:
        prs.save(str(pptx_path))


def get_tmp_project_dir(input_json_path):
    """Derive temp sdpm/{project_name}/ from input JSON path."""
    p = Path(input_json_path).resolve().parent
    project_name = p.name
    tmp_dir = Path(tempfile.gettempdir()) / "sdpm" / project_name
    tmp_dir.mkdir(parents=True, exist_ok=True)
    return tmp_dir
