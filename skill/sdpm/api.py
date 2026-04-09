# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""High-level API for sdpm — single entry points for generate, preview, init, code_block.

These functions encapsulate the full workflow that the CLI (pptx_builder.py) performs.
mcp-local and other consumers should call these instead of assembling low-level APIs.
"""

import os
import sys
from datetime import datetime
from pathlib import Path
from typing import Any


def _resolve_template(data: dict, input_path: str | Path | None, templates_dir: Path) -> tuple[Path, bool]:
    """Resolve template path from presentation data.

    Returns (template_path, custom_template) or raises FileNotFoundError.
    """
    if data.get("template"):
        base_dir = Path(input_path).parent if input_path else Path(".")
        template = base_dir / data["template"]
        if template.exists():
            return template, True
        name = data["template"]
        named = templates_dir / (name if name.endswith(".pptx") else name + ".pptx")
        if named.exists():
            return named, True
    raise FileNotFoundError("No template specified. Set \"template\" in presentation JSON.")


def _get_output_base_dir() -> Path:
    """Get output base directory from config, with WSL fallback."""
    env_dir = os.environ.get("SDPM_OUTPUT_DIR")
    if env_dir:
        return Path(env_dir)
    try:
        from sdpm.config import get_output_dir
        return get_output_dir()
    except Exception:
        pass
    from sdpm.preview import _is_wsl
    if _is_wsl():
        import subprocess
        try:
            result = subprocess.run(  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
                ["/mnt/c/Windows/System32/WindowsPowerShell/v1.0/powershell.exe", "-Command",
                 "[Environment]::GetFolderPath('MyDocuments')"],
                capture_output=True, timeout=10)
            win_path = result.stdout.decode("cp932", errors="replace").strip()
            if win_path:
                wsl = subprocess.run(["wslpath", win_path], capture_output=True, text=True)  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
                if wsl.returncode == 0:
                    return Path(wsl.stdout.strip()) / "SDPM-Presentations"
        except Exception:
            pass
    return Path.home() / "Documents" / "SDPM-Presentations"


def init(
    name: str,
    template: str | Path | None = None,
    output_dir: str | Path | None = None,
) -> dict[str, Any]:
    """Initialize a presentation workspace.

    Creates output directory with presentation.json and specs/.

    Args:
        name: Presentation name (used in directory name).
        template: Template name or path. If provided, extracts fonts.
        output_dir: Explicit output directory. Auto-generated if None.

    Returns:
        Dict with output_dir, json_path, template, fonts, workspace.
    """
    from sdpm.analyzer import extract_fonts
    from sdpm.utils.io import write_json

    if output_dir:
        out_dir = Path(output_dir).expanduser()
    else:
        ts = datetime.now().strftime("%Y%m%d-%H%M")
        dir_name = f"{ts}-{name}" if name else ts
        out_dir = _get_output_base_dir() / dir_name
    out_dir.mkdir(parents=True, exist_ok=True)

    pres_data: dict[str, Any] = {"fonts": {"fullwidth": None, "halfwidth": None}, "slides": []}

    if template:
        templates_dir = Path(__file__).parent.parent / "templates"  # skill/templates/
        template_src = Path(template).expanduser()
        if not template_src.exists():
            candidate = templates_dir / (str(template) if str(template).endswith(".pptx") else f"{template}.pptx")
            if candidate.exists():
                template_src = candidate
        if template_src.exists():
            template_src = template_src.resolve()
            pres_data["template"] = template_src.name
            pres_data["fonts"] = extract_fonts(template_src)

    json_path = out_dir / "presentation.json"
    write_json(json_path, pres_data, suffix="\n")

    specs_dir = out_dir / "specs"
    specs_dir.mkdir(exist_ok=True)
    spec_files = ("brief.md", "outline.md")
    for spec_name in spec_files:
        (specs_dir / spec_name).touch()

    return {
        "output_dir": str(out_dir),
        "json_path": str(json_path),
        "template": pres_data.get("template", ""),
        "fonts": pres_data.get("fonts", {}),
        "workspace": ["presentation.json"] + [f"specs/{s}" for s in spec_files],
    }


def generate(
    json_path: str | Path,
    template: str | Path | None = None,
    output_path: str | Path | None = None,
    no_autofit: bool = False,
    no_preview: bool = False,
) -> dict[str, Any]:
    """Generate PPTX from JSON with full post-processing.

    Includes: template resolution, icon validation, build,
    autofit refresh, unlock_height_constraints, check_layout_imbalance.

    Args:
        json_path: Path to the slides JSON file.
        template: Template override (name or path). Uses JSON's "template" if None.
        output_path: Output .pptx path. Auto-generated if None.
        no_autofit: Skip autofit refresh.
        no_preview: Skip PDF generation during autofit.

    Returns:
        Dict with output_path, slide_count, slides summary, warnings.
    """
    from sdpm.builder import PPTXBuilder, resolve_override, validate_icons_in_json
    from sdpm.preview import (
        check_layout_imbalance_data,
        get_tmp_project_dir,
        refresh_autofit,
        unlock_height_constraints,
    )
    from sdpm.utils.io import read_json

    input_path = Path(json_path)
    if not input_path.exists():
        raise FileNotFoundError(f"Slides JSON not found: {json_path}")

    data = read_json(input_path)
    templates_dir = Path(__file__).parent.parent / "templates"  # skill/templates/
    warnings: list[str] = []

    # Resolve template
    if template:
        p = Path(template)
        if p.exists():
            template_file, custom = p, True
        else:
            named = templates_dir / (str(template) if str(template).endswith(".pptx") else f"{template}.pptx")
            if named.exists():
                template_file, custom = named, True
            else:
                raise FileNotFoundError(f"Template not found: {template}")
    else:
        template_file, custom = _resolve_template(data, str(input_path), templates_dir)

    # Auto-fill fonts / defaultTextColor from template when missing
    from sdpm.analyzer import extract_fonts as _extract_fonts

    fonts = data.get("fonts")
    if not fonts or not fonts.get("fullwidth"):
        fonts = _extract_fonts(template_file)
        warnings.append("fonts auto-detected from template")

    dtc = data.get("defaultTextColor")
    if not dtc:
        _, is_dark = PPTXBuilder._extract_theme_colors(template_file)
        dtc = "#FFFFFF" if is_dark else "#333333"
        warnings.append(f"defaultTextColor auto-set to {dtc}")

    # Validate icons
    missing = validate_icons_in_json(data)
    if missing:
        raise ValueError(f"Missing assets ({len(missing)}): {', '.join(sorted(missing)[:10])}")

    # Build
    builder = PPTXBuilder(
        template_file, custom_template=custom,
        fonts=fonts, base_dir=input_path.parent,
        default_text_color=dtc,
    )
    slides = data.get("slides", [])
    id_map = {}
    for s in slides:
        if "id" in s:
            id_map[s["id"]] = s
    for s in slides:
        builder.add_slide(resolve_override(s, id_map))

    # Output path
    if not output_path:
        ts = datetime.now().strftime("%Y%m%d-%H%M%S")
        p = input_path.with_suffix(".pptx")
        out = p.with_stem(f"{p.stem}_{ts}")
    else:
        out = Path(output_path)

    builder.save(out)

    # Post-processing
    pdf_path = None
    if not no_autofit:
        pptx_resolved = out.resolve()
        pre_autofit_pdf = None
        if not no_preview:
            tmp_project = get_tmp_project_dir(str(input_path))
            preview_dir = tmp_project / "preview"
            preview_dir.mkdir(parents=True, exist_ok=True)
            pre_autofit_pdf = preview_dir / "slides.pdf"
        refresh_autofit(pptx_resolved, pdf_path=None, pre_autofit_pdf=pre_autofit_pdf)
        pdf_path = pre_autofit_pdf
        unlock_height_constraints(pptx_resolved)

    imbalance = check_layout_imbalance_data(out, slides)
    if imbalance:
        for a in imbalance:
            warnings.append(f"page{a['slide']:02d} ({a['layout']}) offset: {a['offset']} ({a['direction']})")

    # Summary
    summary = []
    for i, s in enumerate(slides, 1):
        title = s.get("title", "(no title)")
        if isinstance(title, dict):
            title = title.get("text", "(no title)")
        summary.append(f"page{i:02d} - {title}")

    return {
        "output_path": str(out),
        "slide_count": len(slides),
        "slides": summary,
        "warnings": warnings,
        "pdf_path": str(pdf_path) if pdf_path and pdf_path.exists() else None,
    }


def preview(
    pptx_path: str | Path,
    pages: list[int] | None = None,
    output_dir: str | Path | None = None,
    grid: bool = False,
    pdf_path: str | Path | None = None,
) -> dict[str, Any]:
    """Export PPTX slides as PNG images.

    Uses export_pdf → pdftoppm pipeline on macOS/Linux/WSL.
    Uses COM → PNG direct export on Windows native.

    Args:
        pptx_path: Path to .pptx file.
        pages: Page numbers to export. None for all.
        output_dir: Output directory. Auto-generated if None.
        grid: Add grid overlay to PNGs.
        pdf_path: Pre-generated PDF to skip export step.

    Returns:
        Dict with preview_dir and files list.
    """
    import glob
    import re
    import subprocess

    from pptx import Presentation

    from sdpm.preview import _is_wsl, export_pdf

    path = Path(pptx_path).resolve()
    if not path.exists():
        raise FileNotFoundError(f"PPTX not found: {pptx_path}")

    if output_dir:
        out_dir = Path(output_dir)
    elif _is_wsl() or sys.platform == "win32":
        out_dir = path.parent / "preview"
    else:
        out_dir = Path("/tmp/pptx-preview")
    out_dir.mkdir(parents=True, exist_ok=True)

    pages_set = set(pages) if pages else None

    if sys.platform == "win32":
        generated = _preview_win32(path, out_dir, pages_set)
    else:
        # PDF + pdftoppm pipeline
        pdf = Path(pdf_path) if pdf_path and Path(pdf_path).exists() else out_dir / "slides.pdf"
        if not pdf.exists():
            if not export_pdf(path, pdf):
                raise RuntimeError("PDF export failed. Is a presentation app installed?")

        cmd = ["pdftoppm", "-png", "-scale-to", "1280", str(pdf), str(out_dir / "page")]
        result = subprocess.run(cmd, capture_output=True, text=True)  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
        if result.returncode != 0:
            raise RuntimeError(f"PNG conversion failed. Is poppler (pdftoppm) installed? {result.stderr}")

        # Rename with slide titles
        prs = Presentation(str(path))
        titles = _extract_slide_titles(prs)

        generated = []
        for png in sorted(glob.glob(str(out_dir / "page-*.png"))):
            match = re.match(r'page-(\d+)\.png', Path(png).name)
            if match:
                num = int(match.group(1))
                if pages_set and num not in pages_set:
                    Path(png).unlink()
                    continue
                new_name = f"page{num:02d}-{titles.get(num, 'notitle')}.png"
                new_path = out_dir / new_name
                Path(png).rename(new_path)
                generated.append(str(new_path))

        # Cleanup PDF only if we generated it
        if not pdf_path:
            pdf.unlink(missing_ok=True)

    if grid:
        _apply_grid_overlay(generated)

    return {"preview_dir": str(out_dir), "files": generated}


def _preview_win32(pptx_path: Path, output_dir: Path, pages_set: set[int] | None) -> list[str]:
    """Windows native PNG export via PowerShell COM."""
    import glob
    import shutil
    import subprocess

    from pptx import Presentation

    png_dir = output_dir / "ppt_png"
    png_dir.mkdir(parents=True, exist_ok=True)
    ps_cmd = (
        f"$app = New-Object -ComObject PowerPoint.Application; "
        f"$prs = $app.Presentations.Open('{pptx_path}', "
        f"[Microsoft.Office.Core.MsoTriState]::msoTrue, "
        f"[Microsoft.Office.Core.MsoTriState]::msoFalse, "
        f"[Microsoft.Office.Core.MsoTriState]::msoFalse); "
        f"$prs.SaveAs('{png_dir}', 18); "
        f"$prs.Close(); $app.Quit()"
    )
    result = subprocess.run(["powershell.exe", "-Command", ps_cmd], capture_output=True, timeout=60)  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
    if result.returncode != 0:
        stderr = result.stderr.decode("cp932", errors="replace") if result.stderr else ""
        raise RuntimeError(f"PNG export failed: {stderr}")

    prs = Presentation(str(pptx_path))
    titles = _extract_slide_titles(prs)

    png_files = sorted(glob.glob(str(png_dir / "**" / "*.png"), recursive=True))
    if not png_files:
        png_files = sorted(glob.glob(str(png_dir / "**" / "*.PNG"), recursive=True))

    generated = []
    for idx, png in enumerate(png_files, 1):
        if pages_set and idx not in pages_set:
            Path(png).unlink()
            continue
        new_name = f"page{idx:02d}-{titles.get(idx, 'notitle')}.png"
        new_path = output_dir / new_name
        Path(png).rename(new_path)
        generated.append(str(new_path))

    shutil.rmtree(png_dir, ignore_errors=True)
    return generated


def _extract_slide_titles(prs) -> dict[int, str]:
    """Extract sanitized slide titles from a Presentation object."""
    import re
    titles = {}
    for i, slide in enumerate(prs.slides, 1):
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text.strip().replace("\n", " ")[:30]
        title = re.sub(r'[\\/:*?"<>|]', '', title)
        titles[i] = title or "notitle"
    return titles


def _apply_grid_overlay(png_paths: list[str]) -> None:
    """Add coordinate grid overlay to PNG files."""
    from PIL import Image, ImageDraw, ImageFont

    color = (255, 0, 0, 128)
    try:
        font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 24)
    except Exception:
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
        except Exception:
            font = ImageFont.load_default()

    for png_path in png_paths:
        img = Image.open(png_path).convert("RGBA")
        w, h = img.size
        overlay = Image.new("RGBA", (w, h), (0, 0, 0, 0))
        draw = ImageDraw.Draw(overlay)
        for pct in range(5, 100, 5):
            x, y = int(w * pct / 100), int(h * pct / 100)
            px_x, px_y = int(1920 * pct / 100), int(1080 * pct / 100)
            draw.line([(x, 0), (x, h)], fill=color, width=1)
            draw.line([(0, y), (w, y)], fill=color, width=1)
            if pct % 10 == 0:
                draw.text((x + 4, 4), f"{px_x}px ({pct}%)", fill=color, font=font)
                draw.text((4, y + 4), f"{px_y}px ({pct}%)", fill=color, font=font)
        Image.alpha_composite(img, overlay).convert("RGB").save(png_path)


def code_block(
    code: str,
    language: str = "python",
    theme: str = "dark",
    x: int = 0,
    y: int = 0,
    width: int = 800,
    height: int = 300,
    font_size: int = 12,
    show_label: bool = True,
) -> list[dict[str, Any]]:
    """Generate slide elements for a syntax-highlighted code block.

    Args:
        code: Source code text.
        language: Programming language for highlighting.
        theme: Color theme ("dark" or "light").
        x, y, width, height: Position and size in pixels.
        font_size: Code font size in pt.
        show_label: Show language label bar.

    Returns:
        List of element dicts for slide JSON.
    """
    from sdpm.builder.constants import CODE_COLORS
    from sdpm.utils.text import highlight_code

    colors = CODE_COLORS.get(theme, CODE_COLORS["dark"])
    bg = colors["background"]
    inverse_theme = "light" if theme == "dark" else "dark"
    inverse_bg = CODE_COLORS[inverse_theme]["background"]
    label_fg = "000000" if theme == "dark" else "FFFFFF"
    label_height = 22

    label_map = {"typescript": "TypeScript", "javascript": "JavaScript", "csharp": "C#", "cpp": "C++"}
    label_text = label_map.get(language, language.capitalize())

    elements: list[dict[str, Any]] = []
    if show_label:
        elements.append({
            "type": "textbox",
            "x": x, "y": y, "width": width, "height": label_height,
            "fontSize": 8, "align": "left",
            "fill": inverse_bg,
            "text": f"{{{{#{label_fg}:{label_text}}}}}",
            "marginLeft": 50000, "marginTop": 0, "marginRight": 0, "marginBottom": 0,
            "autoWidth": True,
        })
        code_y = y + label_height
        code_height = height - label_height
    else:
        code_y = y
        code_height = height

    spans = highlight_code(code, language, theme)
    elements.append({
        "type": "textbox",
        "x": x, "y": code_y, "width": width, "height": code_height,
        "fontSize": font_size, "align": "left",
        "fill": bg,
        "text": spans,
    })

    return elements
