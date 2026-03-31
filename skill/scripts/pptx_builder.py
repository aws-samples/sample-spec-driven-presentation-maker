#!/usr/bin/env python3
# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
"""PPTX Builder CLI - backward compatible entry point.

All core logic lives in sdpm package.
This file provides the CLI interface only.
"""
import sys
from pathlib import Path

# Ensure sdpm package is importable
sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import argparse
import json
import re
from datetime import datetime

from pptx import Presentation

from sdpm.assets import (  # noqa: F401
    ICON_DIR,
    ICON_LOCAL_DIR,
    check_asset_exists,
    check_icon_exists,
    print_search_results,
    resolve_asset_path,
    resolve_icon_path,
    search_assets,
)
from sdpm.builder import (
    PPTXBuilder,
    resolve_override,
    validate_icons_in_json,
)
from sdpm.diff import (
    _diff_value,
    _elem_id,
    align_slides,
    load_slides_json_or_pptx,
    match_elements,
)
from sdpm.layout import (
    _layout_collect,
    _layout_route_connections,
    _layout_scale,
    _layout_translate,
    box_to_elements,
)
from sdpm.preview import (
    _is_wsl,
    check_layout_imbalance,
    get_tmp_project_dir,
    refresh_autofit,
    unlock_height_constraints,
)
from sdpm.reference import _get_description
from sdpm.utils.effects import apply_effects  # noqa: F401
from sdpm.utils.image import apply_image_effects, resolve_image_path  # noqa: F401
from sdpm.utils.io import read_json, write_json
from sdpm.utils.svg import (  # noqa: F401
    _recolor_svg,
    add_svg_to_slide,
    generate_qr_svg,
    get_svg_dimensions,
)

# Re-export for backward compatibility (scripts that import from here)
from sdpm.utils.text import normalize_spacing, parse_styled_text  # noqa: F401


def _resolve_template(data, input_path):
    """Resolve template path: presentation.json "template" → templates/default-dark.pptx → repo fallback."""
    templates_dir = Path(__file__).parent.parent / "templates"

    if data.get("template"):
        base_dir = Path(input_path).parent if input_path and input_path != "-" else Path(".")
        template = base_dir / data["template"]
        if template.exists():
            return template, True
        # Resolve by name from repo templates/
        name = data["template"]
        named = templates_dir / (name if name.endswith(".pptx") else name + ".pptx")
        if named.exists():
            return named, True

    default = templates_dir / "default-dark.pptx"
    if default.exists():
        return default, False

    print(f"Error: No template found. Place a template in {templates_dir}/ or set \"template\" in presentation JSON", file=sys.stderr)
    sys.exit(1)


def cmd_generate(args):
    """Generate PPTX from JSON."""
    # Read input data first
    if args.input and args.input != "-":
        data = read_json(Path(args.input))
    else:
        data = json.load(sys.stdin)

    template, custom_template = _resolve_template(
        data,
        args.input if hasattr(args, 'input') else None
    )

    if not template.exists():
        print(f"Error: Template not found: {template}", file=sys.stderr)
        sys.exit(1)

    missing_icons = validate_icons_in_json(data)
    if missing_icons:
        print(f"Error: Missing assets ({len(missing_icons)}):", file=sys.stderr)
        for icon in sorted(missing_icons):
            print(f"  - {icon}", file=sys.stderr)
        print("", file=sys.stderr)
        print("Run the following command to download assets:", file=sys.stderr)
        print("  python3 scripts/download_aws_icons.py", file=sys.stderr)
        print("  python3 scripts/download_material_icons.py", file=sys.stderr)
        sys.exit(1)

    base_dir = Path(args.input).parent if args.input and args.input != "-" else Path(".")
    builder = PPTXBuilder(template, custom_template=custom_template,
                          fonts=data.get("fonts"), base_dir=base_dir,
                          keep_empty_placeholders=getattr(args, 'keep_empty_placeholders', False),
                          default_text_color=data.get("defaultTextColor"))
    slides = data.get("slides", [])

    id_map = {}
    for slide_def in slides:
        if "id" in slide_def:
            sid = slide_def["id"]
            if sid in id_map:
                print(f"Error: Duplicate slide id: {sid}", file=sys.stderr)
                sys.exit(1)
            id_map[sid] = slide_def

    for slide_def in slides:
        resolved = resolve_override(slide_def, id_map)
        builder.add_slide(resolved)
    ts = datetime.now().strftime("%Y%m%d-%H%M%S")
    p = Path(args.output)
    output_path = p.with_stem(f"{p.stem}_{ts}")
    builder.save(output_path)

    print(f"Generated: {output_path.resolve()}")
    for i, slide_def in enumerate(slides, 1):
        title = (slide_def.get("placeholders") or {}).get("0", "(no title)")
        if isinstance(title, dict):
            title = title.get("text", "(no title)")
        print(f"page{i:02d} - {title}")

    pdf_path = None
    if not args.no_autofit:
        pptx_path = output_path.resolve()
        if not args.no_preview:
            preview_dir = Path("/tmp/pptx-preview")
            preview_dir.mkdir(parents=True, exist_ok=True)
            pdf_path = preview_dir / "slides.pdf"
        refresh_autofit(pptx_path, pdf_path=pdf_path)
        unlock_height_constraints(pptx_path)
    check_layout_imbalance(output_path, slides)

    tmp_project = None
    if args.input and args.input != "-":
        tmp_project = get_tmp_project_dir(args.input)

    if not args.no_preview:
        preview_dir = tmp_project / 'preview' if tmp_project else Path("/tmp/pptx-preview")
        preview_dir.mkdir(parents=True, exist_ok=True)
        preview_args = argparse.Namespace(
            input=str(output_path),
            pages=None,
            no_grid=True,
            _pdf_path=str(pdf_path) if pdf_path and pdf_path.exists() else None,
            _output_dir=str(preview_dir),
        )
        cmd_preview(preview_args)


def _export_png_win32(pptx_path, output_dir):
    """Export PPTX to PNG via PowerShell COM (Windows native)."""
    import subprocess
    png_dir = output_dir / "ppt_png"
    png_dir.mkdir(parents=True, exist_ok=True)
    ps_cmd = (
        f"$app = New-Object -ComObject PowerPoint.Application; "
        f"$prs = $app.Presentations.Open('{pptx_path}', "
        f"[Microsoft.Office.Core.MsoTriState]::msoTrue, "
        f"[Microsoft.Office.Core.MsoTriState]::msoFalse, "
        f"[Microsoft.Office.Core.MsoTriState]::msoFalse); "
        f"$prs.SaveAs('{png_dir}', 18); "
        f"$prs.Close(); $app.Quit()"
    )
    result = subprocess.run(["powershell.exe", "-Command", ps_cmd], capture_output=True, timeout=60)  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
    stderr = result.stderr.decode("cp932", errors="replace") if result.stderr else ""
    if result.returncode != 0:
        print(f"Error: PNG export failed: {stderr}", file=sys.stderr)
        sys.exit(1)
    return png_dir
def cmd_preview(args):
    """Export PPTX slides as PNG images."""
    import glob
    import subprocess
    
    pptx_path = Path(args.input).resolve()
    
    if _is_wsl() or sys.platform == "win32":
        output_dir = pptx_path.parent / "preview"
    elif hasattr(args, '_output_dir') and args._output_dir:
        output_dir = Path(args._output_dir)
    else:
        output_dir = Path("/tmp/pptx-preview")
    output_dir.mkdir(parents=True, exist_ok=True)
    
    if not pptx_path.exists():
        print(f"Error: File not found: {pptx_path}", file=sys.stderr)
        sys.exit(1)
    
    if sys.platform == "win32":
        # Windows native: direct PNG export via COM
        png_dir = _export_png_win32(pptx_path, output_dir)
        
        # Parse pages to export
        pages_to_export = None
        if args.pages:
            pages_to_export = set(int(p.strip()) for p in args.pages.split(","))
        
        # Collect PNGs from PowerPoint output
        png_files = sorted(glob.glob(str(png_dir / "**" / "*.png"), recursive=True))
        if not png_files:
            png_files = sorted(glob.glob(str(png_dir / "**" / "*.PNG"), recursive=True))
        
        # Rename with slide titles
        prs = Presentation(str(pptx_path))
        titles = {}
        for i, slide in enumerate(prs.slides, 1):
            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text.strip().replace("\n", " ")[:30]
            title = re.sub(r'[\\/:*?"<>|]', '', title)
            titles[i] = title or "notitle"
        
        generated = []
        for idx, png in enumerate(png_files, 1):
            if pages_to_export and idx not in pages_to_export:
                Path(png).unlink()
                continue
            new_name = f"page{idx:02d}-{titles.get(idx, 'notitle')}.png"
            new_path = output_dir / new_name
            Path(png).rename(new_path)
            generated.append(new_path)
        
        # Cleanup temp PNG dir
        import shutil
        shutil.rmtree(png_dir, ignore_errors=True)
    else:
        # macOS / WSL: PDF + pdftoppm pipeline
        pdf_path = output_dir / "slides.pdf"
        pre_generated_pdf = getattr(args, '_pdf_path', None)
        
        if pre_generated_pdf and Path(pre_generated_pdf).exists():
            # PDF already generated by _refresh_autofit — skip PowerPoint
            import shutil
            if str(Path(pre_generated_pdf).resolve()) != str(pdf_path.resolve()):
                shutil.copy2(pre_generated_pdf, pdf_path)
        else:
            if _is_wsl():
                win_pptx = subprocess.run(["wslpath", "-w", str(pptx_path)], capture_output=True, text=True).stdout.strip()  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
                win_pdf = subprocess.run(["wslpath", "-w", str(pdf_path)], capture_output=True, text=True).stdout.strip()  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
                ps_cmd = (
                    f"$app = New-Object -ComObject PowerPoint.Application; "
                    f"$prs = $app.Presentations.Open('{win_pptx}', "
                    f"[Microsoft.Office.Core.MsoTriState]::msoTrue, "
                    f"[Microsoft.Office.Core.MsoTriState]::msoFalse, "
                    f"[Microsoft.Office.Core.MsoTriState]::msoFalse); "
                    f"$prs.SaveAs('{win_pdf}', 32); "
                    f"$prs.Close(); $app.Quit()"
                )
                result = subprocess.run(
                    ["/mnt/c/Windows/System32/WindowsPowerShell/v1.0/powershell.exe", "-Command", ps_cmd],
                    capture_output=True, timeout=60)  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
                result = subprocess.CompletedProcess(
                    result.args, result.returncode,
                    stdout=result.stdout.decode("cp932", errors="replace") if result.stdout else "",
                    stderr=result.stderr.decode("cp932", errors="replace") if result.stderr else ""
                )
            else:
                from sdpm.preview import (
                    _mac_open_pptx_background,
                    _mac_restore_pptx_focus_async,
                )
                pptx_name = Path(pptx_path).name
                restore_info = _mac_open_pptx_background(pptx_path)
                _mac_restore_pptx_focus_async(restore_info)
                import time
                time.sleep(2)
                script = f'''
                tell application "Microsoft PowerPoint"
                    set theDoc to presentation "{pptx_name}"
                    set outPath to (POSIX file "{pdf_path}") as text
                    save theDoc in outPath as save as PDF
                    close theDoc saving no
                end tell
                '''
                result = subprocess.run(["osascript", "-e", script], capture_output=True, text=True)  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
            
            if result.returncode != 0 or not pdf_path.exists():
                print(f"Error: PDF export failed: {result.stderr}", file=sys.stderr)
                sys.exit(1)
        
        # Parse pages to export
        pages_to_export = None
        if args.pages:
            pages_to_export = set(int(p.strip()) for p in args.pages.split(","))
        
        # PDF -> PNG via pdftoppm (max 1280px wide for compact preview)
        cmd = ["pdftoppm", "-png", "-scale-to", "1280", str(pdf_path), str(output_dir / "page")]
        result = subprocess.run(cmd, capture_output=True, text=True)  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
        if result.returncode != 0:
            print(f"Error: PNG conversion failed: {result.stderr}", file=sys.stderr)
            sys.exit(1)
        
        # Rename with slide titles
        prs = Presentation(str(pptx_path))
        titles = {}
        for i, slide in enumerate(prs.slides, 1):
            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text.strip().replace("\n", " ")[:30]
            title = re.sub(r'[\\/:*?"<>|]', '', title)
            titles[i] = title or "notitle"
        
        generated = []
        for png in sorted(glob.glob(str(output_dir / "page-*.png"))):
            basename = Path(png).name
            match = re.match(r'page-(\d+)\.png', basename)
            if match:
                num = int(match.group(1))
                if pages_to_export and num not in pages_to_export:
                    Path(png).unlink()
                    continue
                new_name = f"page{num:02d}-{titles.get(num, 'notitle')}.png"
                new_path = output_dir / new_name
                Path(png).rename(new_path)
                generated.append(new_path)
        
        # Cleanup PDF
        pdf_path.unlink()
    
    # Add grid overlay unless disabled
    if not args.no_grid:
        from PIL import Image, ImageDraw, ImageFont
        color = (255, 0, 0, 128)
        try:
            font = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 24)
        except Exception:
            try:
                font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 24)
            except Exception:
                font = ImageFont.load_default()
        
        for png_path in generated:
            img = Image.open(png_path).convert("RGBA")
            w, h = img.size
            overlay = Image.new("RGBA", (w, h), (0, 0, 0, 0))
            draw = ImageDraw.Draw(overlay)
            for pct in range(5, 100, 5):
                x, y = int(w * pct / 100), int(h * pct / 100)
                px_x, px_y = int(1920 * pct / 100), int(1080 * pct / 100)
                draw.line([(x, 0), (x, h)], fill=color, width=1)
                draw.line([(0, y), (w, y)], fill=color, width=1)
                if pct % 10 == 0:
                    draw.text((x + 4, 4), f"{px_x}px ({pct}%)", fill=color, font=font)
                    draw.text((4, y + 4), f"{px_y}px ({pct}%)", fill=color, font=font)
            Image.alpha_composite(img, overlay).convert("RGB").save(png_path)
    
    for path in generated:
        print(f"Generated: {path}")
    if generated:
        print(f"Preview: {generated[0].parent}")
    if generated:
        print(f"Preview: {generated[0].parent}")





def cmd_search_assets(args):
    """Search assets (icons, images, etc.) by keywords."""
    results = search_assets(
        query=args.query,
        limit=args.limit,
        source_filter=args.source,
        type_filter=args.type,
        theme_filter=args.theme,
    )
    print_search_results(results, limit=args.limit)


# Backward-compatible alias
cmd_icon_search = cmd_search_assets


def cmd_list_asset_sources(args):
    """List available asset sources."""
    from sdpm.assets import list_sources
    sources = list_sources()
    if not sources:
        print("No asset sources found.", file=sys.stderr)
        return
    for s in sources:
        desc = f"  {s['description']}" if s["description"] else ""
        print(f"  {s['source']:<20} {s['count']:>6} assets{desc}")


def cmd_list_templates(args):
    """List available PPTX templates."""
    templates_dir = Path(__file__).parent.parent / "templates"
    if not templates_dir.exists():
        print("No templates directory found.", file=sys.stderr)
        return
    templates = sorted(templates_dir.glob("*.pptx"))
    if not templates:
        print("No templates found.", file=sys.stderr)
        return
    for t in templates:
        print(f"  {t.stem}")


def _list_or_show(category_dir, names, category_label):
    """List or show documents in a category directory.

    Supports two layouts:
    - Categorized: subdirectories as categories, names use 'category/name' format
    - Flat: files directly in category_dir, names use plain 'name' format

    Args:
        category_dir: Path to the category directory.
        names: List of names to show. Empty list shows the listing.
        category_label: Display label (e.g. "Examples").
    """
    if not category_dir.exists():
        print(f"Directory not found: {category_dir}", file=sys.stderr)
        return

    # Discover subdirectories as categories
    subdirs = sorted(d for d in category_dir.iterdir() if d.is_dir())
    flat = not subdirs
    root_files = _collect_files(category_dir) if flat else {}

    if names:
        for name in names:
            if flat or '/' not in name:
                # Flat mode: look up directly in category_dir
                files = root_files if flat else {}
                if not flat:
                    # Categorized mode but no slash — try to find in subdirs
                    for sd in subdirs:
                        files = _collect_files(sd)
                        if name in files:
                            print(f"# Format: category/name (e.g. {sd.name}/{name})", file=sys.stderr)
                            break
                    else:
                        print(f"# Format: category/name (e.g. pattern/{name})", file=sys.stderr)
                    continue
                doc = files.get(name)
                if not doc:
                    print(f"# Not found: {name}", file=sys.stderr)
                    print(f"# Available: {', '.join(sorted(files.keys()))}", file=sys.stderr)
                    continue
            else:
                cat, stem = name.split('/', 1)
                sub = category_dir / cat
                if not sub.is_dir():
                    cats = [d.name for d in subdirs]
                    print(f"# Category not found: {cat}", file=sys.stderr)
                    print(f"# Available: {', '.join(cats)}", file=sys.stderr)
                    continue
                files = _collect_files(sub)
                doc = files.get(stem)
                if not doc:
                    print(f"# Not found: {name}", file=sys.stderr)
                    print(f"# Available in {cat}/: {', '.join(sorted(files.keys()))}", file=sys.stderr)
                    continue
            if doc.suffix == '.md':
                text = doc.read_text(encoding="utf-8")
                # Strip YAML frontmatter
                lines = text.splitlines(True)
                if lines and lines[0].strip() == '---':
                    for i, line in enumerate(lines[1:], 1):
                        if line.strip() == '---':
                            text = ''.join(lines[i + 1:]).lstrip('\n')
                            break
                print(text)
            print()
    else:
        print(f"# {category_label}")
        if flat:
            for stem in sorted(root_files.keys()):
                f = root_files[stem]
                desc = _get_description(f) or stem
                print(f"  {stem:<36} {desc}")
        else:
            for sd in subdirs:
                files = _collect_files(sd)
                if not files:
                    continue
                print(f"\n## {sd.name.title()}")
                for stem in sorted(files.keys()):
                    f = files[stem]
                    desc = _get_description(f) or stem
                    print(f"  {sd.name}/{stem:<28} {desc}")


def _collect_files(directory):
    """Collect md/pptx files in directory. md takes priority over pptx for same stem."""
    files = {}
    for f in sorted(directory.glob("*.pptx")):
        files[f.stem] = f
    for f in sorted(directory.glob("*.md")):
        files[f.stem] = f
    return files


def cmd_search_patterns(args):
    """Search patterns by keywords."""
    from sdpm.reference import search_patterns
    results = search_patterns(args.query, limit=args.limit)
    if not results:
        print("No matches found.")
        return
    for r in results:
        page = f"/{r['page']}" if r.get('page') else ""
        print(f"  {r['path']}{page}  {r['description']}")


def cmd_examples(args):
    """List or show design examples (components/patterns/styles)."""
    from sdpm.reference import get_pptx_notes, list_pptx_descriptions

    examples_dir = Path(__file__).parent.parent / "references" / "examples"
    if not examples_dir.exists():
        print("Directory not found: references/examples", file=sys.stderr)
        return

    names = args.names
    if not names:
        print("Usage: examples <category> or <category/name>", file=sys.stderr)
        return

    for name in names:
        parts = name.split("/", 1)
        base = parts[0]
        sub = parts[1] if len(parts) > 1 else None

        # styles/ directory — always show full content
        if base == "styles":
            styles_dir = examples_dir / "styles"
            if not styles_dir.exists():
                print("# Not found: styles/", file=sys.stderr)
                continue
            if sub is None:
                # List styles
                for f in sorted(styles_dir.iterdir()):
                    if f.suffix == '.html' and not f.name.startswith('.'):
                        from sdpm.reference import _get_description
                        desc = _get_description(f)
                        print(f"  styles/{f.stem}  {desc}")
                # Open index in browser unless --no-browse
                if not args.no_browse:
                    import webbrowser
                    from sdpm.reference import generate_styles_index
                    index_path = generate_styles_index(styles_dir)
                    if index_path.exists():
                        webbrowser.open(index_path.as_uri())
            else:
                print("# Use 'init --style {name}' to copy a style, then read specs/art-direction.html", file=sys.stderr)
            continue

        # pptx files (components, patterns) — directory-like behavior
        pptx_path = examples_dir / f"{base}.pptx"
        if not pptx_path.is_file():
            print(f"# Not found: {base}", file=sys.stderr)
            cats = []
            for f in sorted(examples_dir.iterdir()):
                if f.suffix == ".pptx":
                    cats.append(f.stem)
                elif f.is_dir() and not f.name.startswith('.'):
                    cats.append(f"{f.name}/")
            print(f"# Available: {', '.join(cats)}", file=sys.stderr)
            continue

        if sub is None:
            # List slides
            for page, desc in list_pptx_descriptions(str(pptx_path)):
                print(f"  {page:>3}  {desc}")
        elif sub == "all":
            # All notes
            for _, notes in get_pptx_notes(str(pptx_path)):
                print(notes)
                print()
        else:
            # Specific pages
            try:
                pages = [int(p.strip()) for p in sub.split(",")]
            except ValueError:
                print(f"# Invalid page: {sub}", file=sys.stderr)
                continue
            results = get_pptx_notes(str(pptx_path), pages=pages)
            if not results:
                print(f"# No notes found for page(s): {sub}", file=sys.stderr)
            for _, notes in results:
                print(notes)
                print()


def cmd_workflows(args):
    """List or show workflow documents."""
    _list_or_show(Path(__file__).parent.parent / "references" / "workflows", args.names, "Workflows")


def cmd_guides(args):
    """List or show guide documents."""
    _list_or_show(Path(__file__).parent.parent / "references" / "guides", args.names, "Guides")


def _get_documents_dir():
    """Get Documents directory."""
    import subprocess
    if _is_wsl():
        try:
            result = subprocess.run(  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
                ["/mnt/c/Windows/System32/WindowsPowerShell/v1.0/powershell.exe", "-Command",
                 "[Environment]::GetFolderPath('MyDocuments')"],
                capture_output=True, timeout=10)
            win_path = result.stdout.decode("cp932", errors="replace").strip()
            if win_path:
                wsl = subprocess.run(["wslpath", win_path], capture_output=True, text=True)  # nosec B603 # nosemgrep: python.lang.security.audit.dangerous-subprocess-use-audit
                if wsl.returncode == 0:
                    return Path(wsl.stdout.strip())
        except Exception:
            pass
    return Path.home() / "Documents"


def cmd_init(args):
    if args.output:
        out_dir = Path(args.output).expanduser()
    else:
        ts = datetime.now().strftime("%Y%m%d-%H%M")
        name = f"{ts}-{args.name}" if args.name else ts
        out_dir = _get_documents_dir() / "pptx-maker" / name
    out_dir.mkdir(parents=True, exist_ok=True)
    pres_data = {"fonts": {"fullwidth": None, "halfwidth": None}, "slides": []}
    json_path = out_dir / "presentation.json"
    write_json(json_path, pres_data, suffix="\n")
    specs_dir = out_dir / "specs"
    specs_dir.mkdir(exist_ok=True)
    spec_files = ["brief.md", "outline.md"]
    if args.style:
        # Copy style HTML as art-direction.html
        styles_dir = Path(__file__).parent.parent / "references" / "examples" / "styles"
        style_name = args.style if args.style.endswith(".html") else args.style + ".html"
        style_src = styles_dir / style_name
        if not style_src.exists():
            print(f"Error: Style not found: {style_src}", file=sys.stderr)
            sys.exit(1)
        import shutil
        shutil.copy2(style_src, specs_dir / "art-direction.html")
        print(f"style:       {style_src.name} → specs/art-direction.html")
    else:
        spec_files.append("art-direction.md")
    for name in spec_files:
        (specs_dir / name).touch()
    print(f"output_json: {json_path}")
    for name in spec_files:
        print(f"specs:       {specs_dir / name}")
def cmd_code_block(args):
    """Generate elements JSON for a syntax-highlighted code block."""
    from sdpm.builder.constants import CODE_COLORS
    from sdpm.utils.text import highlight_code

    if args.input == "-":
        import sys as _sys
        code = _sys.stdin.read()
    else:
        code = Path(args.input).read_text(encoding="utf-8")

    language = args.language or "text"
    theme = args.theme or "dark"
    x = args.x or 0
    y = args.y or 0
    width = args.width or 800
    height = args.height or 300
    label_height = 22
    show_label = not args.no_label

    colors = CODE_COLORS.get(theme, CODE_COLORS["dark"])
    bg = colors["background"]
    inverse_theme = "light" if theme == "dark" else "dark"
    inverse_bg = f"#{CODE_COLORS[inverse_theme]['background'].lstrip('#')}" if isinstance(CODE_COLORS[inverse_theme]['background'], str) else CODE_COLORS[inverse_theme]['background']
    # Resolve string hex colors
    if isinstance(bg, str):
        bg_hex = bg
    else:
        bg_hex = f"#{bg:06x}" if isinstance(bg, int) else bg

    highlighted = highlight_code(code, language, theme)
    label_map = {"typescript": "TypeScript", "javascript": "JavaScript", "csharp": "C#", "cpp": "C++"}
    label_text = label_map.get(language, language.capitalize())

    elements = []
    if show_label:
        elements.append({
            "type": "textbox",
            "x": x, "y": y, "width": width, "height": label_height,
            "fontSize": 8, "align": "left",
            "fill": inverse_bg,
            "text": f"{{{{#{'000000' if theme == 'dark' else 'FFFFFF'}:{label_text}}}}}",
            "marginLeft": 50000, "marginTop": 0, "marginRight": 0, "marginBottom": 0,
            "autoWidth": True,
        })
        code_y = y + label_height
        code_height = height - label_height
    else:
        code_y = y
        code_height = height

    elements.append({
        "type": "textbox",
        "x": x, "y": code_y, "width": width, "height": code_height,
        "fontSize": args.font_size or 12, "align": "left",
        "fill": bg_hex,
        "text": highlighted,
    })

    output = {"elements": elements}
    out_str = json.dumps(output, ensure_ascii=False, indent=2)
    if args.output:
        write_json(Path(args.output), output)
        print(f"Written: {args.output}", file=__import__('sys').stderr)
    else:
        print(out_str)


def cmd_layout(args):
    """Layout engine: compute coordinates from logical structure JSON."""
    if args.input == "-":
        import sys as _sys
        source = _sys.stdin.read()
    else:
        source = Path(args.input).read_text(encoding="utf-8")

    tree = json.loads(source)
    direction = tree.get("direction", "horizontal")
    align = tree.get("align", "center")
    reverse = tree.get("reverse", False)

    def build_root():
        import copy
        root = {"id": "_root", "children": copy.deepcopy(tree.get("children", tree.get("nodes", []))),
                "direction": direction, "align": align}
        if reverse:
            root["reverse"] = True
        return root

    # Pass 1: natural size
    root = build_root()
    _layout_scale(root, direction, align)
    rb = root["_bindings"]
    _, _ = rb[2], rb[3]

    target_w = args.width
    target_h = args.height

    cumulative_sh = 1.0
    cumulative_sv = 1.0
    if target_w or target_h:
        for _ in range(8):
            rb = root["_bindings"]
            sx = target_w / rb[2] if target_w else 1.0
            sy = target_h / rb[3] if target_h else 1.0
            if abs(sx - 1.0) < 0.03 and abs(sy - 1.0) < 0.03:
                break
            cumulative_sh *= sx
            cumulative_sv *= sy
            root = build_root()
            _layout_scale(root, direction, align, cumulative_sh, cumulative_sv)
        args._cumulative_scale = min(cumulative_sh, cumulative_sv)

    rb = root["_bindings"]
    ox = (args.x or 0) - rb[0]
    oy = (args.y or 0) - rb[1]
    _layout_translate(root, ox, oy)

    # Center if still slightly off from target
    rb = root["_bindings"]
    if target_w and abs(rb[2] - target_w) > 5:
        dx = (target_w - rb[2]) // 2
        _layout_translate(root, dx, 0)
    if target_h and abs(rb[3] - target_h) > 5:
        dy = (target_h - rb[3]) // 2
        _layout_translate(root, 0, dy)

    # Collect results
    nodes_out = {}
    groups_out = {}
    rb = root["_bindings"]
    for child in root["children"]:
        _layout_collect(child, nodes_out, groups_out)

    # Route connections
    edges_out = []
    connections = tree.get("connections", [])
    if connections:
        edges_out = _layout_route_connections(connections, nodes_out, groups_out)

    # Build pptx-maker elements array
    elements = []

    # Groups (largest first for correct z-order)
    for gid, g in sorted(groups_out.items(), key=lambda x: -x[1]["width"] * x[1]["height"]):
        gt = g.get("groupType")
        if not gt:
            continue
        elements.append({"type": "arch-group", "groupType": gt, "x": g["x"], "y": g["y"], "width": g["width"], "height": g["height"], "label": g.get("label", gid.rsplit(".", 1)[-1])})

    # Nodes
    is_dark = getattr(args, 'theme', 'dark') == 'dark'
    for nid, n in nodes_out.items():
        if n.get("box"):
            elements.extend(box_to_elements(nid, n, is_dark))
        elif n.get("icon"):
            elements.append({"type": "image", "src": n["icon"], "x": n["x"], "y": n["y"], "width": n["width"], "label": n.get("label", nid.rsplit(".", 1)[-1]), "labelPosition": "bottom"})

    # Edges as elbow connectors
    # Track label positions for overlap avoidance
    placed_labels = []

    for e in edges_out:
        pts = e["points"]
        if len(pts) < 2:
            continue
        sx, sy = pts[0]
        ex, ey = pts[-1]
        el = {"type": "line", "x1": sx, "y1": sy, "x2": ex, "y2": ey, "arrowEnd": "arrow"}
        if len(pts) > 2:
            el["connectorType"] = "elbow"
            dx = ex - sx
            dy = ey - sy
            if len(pts) >= 4 and abs(dx) > 0 and abs(dy) > 0:
                # 4 points: [start, bend1, bend2, end]
                seg1_vertical = abs(pts[0][0] - pts[1][0]) <= abs(pts[0][1] - pts[1][1])
                if seg1_vertical:
                    # V-H-V → elbowStart vertical
                    adj = (pts[1][1] - sy) / dy if dy != 0 else 0.5
                    el["preset"] = "bentConnector3"
                    el["elbowStart"] = "vertical"
                    el["adjustments"] = [max(0.0, min(1.0, adj))]
                else:
                    # H-V-H → bentConnector4 (no flip)
                    adj1 = (pts[1][0] - sx) / dx if dx != 0 else 0.5
                    adj2 = (pts[2][1] - sy) / dy if dy != 0 else 0.5
                    el["preset"] = "bentConnector4"
                    el["adjustments"] = [max(-1.0, min(2.0, adj1)), max(-1.0, min(2.0, adj2))]
            elif dy != 0 or dx != 0:
                el["adjustments"] = [0.5]
        elements.append(el)

        label = e.get("label", "")
        if not label:
            continue

        # Apply user labelOffset if provided
        conn_obj = None
        for c in (tree.get("connections") or []):
            if c.get("from") == e["from"] and c.get("to") == e["to"]:
                conn_obj = c
                break
        user_offset = (conn_obj or {}).get("labelOffset", {})

        # Find longest segment midpoint
        best_mid = None
        best_len = -1
        best_horizontal = True
        for si in range(len(pts) - 1):
            ax, ay = pts[si]
            bx, by = pts[si + 1]
            seg_len = abs(bx - ax) + abs(by - ay)
            if seg_len > best_len:
                best_len = seg_len
                best_mid = ((ax + bx) // 2, (ay + by) // 2)
                best_horizontal = abs(bx - ax) > abs(by - ay)
        mx, my = best_mid or ((pts[0][0] + pts[-1][0]) // 2, (pts[0][1] + pts[-1][1]) // 2)

        tw = max(len(label) * 11, 60)
        th = 30
        arrow_len = abs(ex - sx) + abs(ey - sy)

        # Position: below for horizontal, right for vertical
        if best_horizontal:
            lx, ly = mx - tw // 2, my + 2
            # Short arrow: place label below arrow end
            if arrow_len < tw + 20:
                ly = my + 12
        else:
            lx, ly = mx + 6, my - th // 2

        # Apply user offset
        lx += user_offset.get("x", 0)
        ly += user_offset.get("y", 0)

        # Overlap avoidance: shift if colliding with existing labels
        for px, py, pw, ph in placed_labels:
            if lx < px + pw and lx + tw > px and ly < py + ph and ly + th > py:
                if best_horizontal:
                    ly = py + ph + 2
                else:
                    ly = py + ph + 2
        placed_labels.append((lx, ly, tw, th))

        elements.append({"type": "textbox", "x": lx, "y": ly, "width": tw, "height": th, "fontSize": 9, "align": "center", "verticalAlign": "top", "fill": "#000000", "opacity": 0.7, "line": "none", "marginTop": 0, "marginBottom": 0, "marginLeft": 0, "marginRight": 0, "text": "{{#8FA7C4:" + label + "}}"})

    output = {
        "elements": elements,
        "bbox": {"x": rb[0], "y": rb[1], "width": rb[2], "height": rb[3]},
    }

    # Generate layout feedback
    warnings = []
    if target_w:
        ratio_w = rb[2] / target_w
        if ratio_w < 0.5:
            warnings.append(f"Layout uses only {round(ratio_w*100)}% of target width ({rb[2]}px / {target_w}px). Consider placing top-level groups horizontally.")
        if rb[2] > target_w:
            warnings.append(f"Layout width {rb[2]}px exceeds target {target_w}px. Consider reducing horizontal elements or splitting into multiple rows.")
    if target_h:
        ratio_h = rb[3] / target_h
        if ratio_h < 0.5:
            warnings.append(f"Layout uses only {round(ratio_h*100)}% of target height ({rb[3]}px / {target_h}px). Consider adding vertical spacing or stacking groups vertically.")
        if rb[3] > target_h:
            warnings.append(f"Layout height {rb[3]}px exceeds target {target_h}px. Consider reducing nesting depth or placing groups horizontally.")
    if hasattr(args, '_cumulative_scale'):
        if cumulative_sh < 0.5:
            warnings.append(f"Horizontal spacing compressed to {round(cumulative_sh*100)}%. Consider reducing horizontal elements.")
        if cumulative_sv < 0.5:
            warnings.append(f"Vertical spacing compressed to {round(cumulative_sv*100)}%. Consider reducing vertical stacking.")

    # Check per-group size
    for gid, g in groups_out.items():
        children = g.get("children", [])
        if len(children) >= 3:
            if target_h and g["height"] > (target_h * 0.6):
                warnings.append(f"Group \"{g['label']}\" is tall ({g['height']}px). Consider direction: horizontal for its children.")
            if target_w and g["width"] > (target_w * 0.8):
                warnings.append(f"Group \"{g['label']}\" is wide ({g['width']}px). Consider direction: vertical for its children.")

    # Check label overlaps
    label_rects = []
    for nid, n in nodes_out.items():
        lbl = n.get("label", "")
        if lbl:
            lw = len(lbl) * 8 + 10
            lh = 20
            lx = n["x"] + (n["width"] - lw) / 2
            ly = n["y"] + n["height"]
            label_rects.append((nid, lbl, lx, ly, lw, lh))
    for i in range(len(label_rects)):
        for j in range(i + 1, len(label_rects)):
            _, l1, x1, y1, w1, h1 = label_rects[i]
            _, l2, x2, y2, w2, h2 = label_rects[j]
            gap = 5
            if x1 - gap < x2 + w2 and x1 + w1 + gap > x2 and y1 - gap < y2 + h2 and y1 + h1 + gap > y2:
                warnings.append(f"Labels \"{l1}\" and \"{l2}\" overlap. Increase spacing or shorten labels.")

    # Check edge-node crossings
    margin = 5
    crossing_reported = set()
    for e in edges_out:
        pts = e["points"]
        if len(pts) < 2:
            continue
        src_id, dst_id = e["from"], e["to"]
        edge_key = f"{src_id}→{dst_id}"
        for seg_i in range(len(pts) - 1):
            x1, y1 = pts[seg_i]
            x2, y2 = pts[seg_i + 1]
            seg_min_x, seg_max_x = min(x1, x2), max(x1, x2)
            seg_min_y, seg_max_y = min(y1, y2), max(y1, y2)
            for nid, n in nodes_out.items():
                if nid.endswith(src_id) or nid.endswith(dst_id):
                    continue
                report_key = (edge_key, n.get("label", nid))
                if report_key in crossing_reported:
                    continue
                nx, ny, nw, nh = n["x"], n["y"], n["width"], n["height"]
                if seg_max_x > nx + margin and seg_min_x < nx + nw - margin and seg_max_y > ny + margin and seg_min_y < ny + nh - margin:
                    # Suggest reordering based on connection direction
                    suggest = ""
                    dst_node = nodes_out.get(dst_id) or next((v for k, v in nodes_out.items() if k.endswith(dst_id)), None)
                    if dst_node:
                        src_node = nodes_out.get(src_id) or next((v for k, v in nodes_out.items() if k.endswith(src_id)), None)
                        if src_node:
                            dx = dst_node["x"] - src_node["x"]
                            direction = "rightmost" if dx > 0 else "leftmost"
                            suggest = f' Suggest: place "{src_node.get("label", src_id)}" {direction} in its group, adjacent to "{dst_node.get("label", dst_id)}".'
                    warnings.append(f'Edge {edge_key} crosses node "{n.get("label", nid)}". Reorder nodes so connected elements are adjacent, or group branch targets in the perpendicular direction. Also consider reverse: true on the target group if connections flow opposite to layout direction.{suggest}')
                    crossing_reported.add(report_key)

    # Structure suggestions: sibling size imbalance
    all_items = {}
    all_items.update(groups_out)
    all_items.update(nodes_out)
    for gid, g in groups_out.items():
        child_ids = g.get("children", [])
        if len(child_ids) < 2:
            continue
        has_group_child = any(cid in groups_out for cid in child_ids)
        if not has_group_child:
            continue
        child_bboxes = []
        for cid in child_ids:
            c = all_items.get(cid)
            if c:
                child_bboxes.append((cid, c))
        if len(child_bboxes) < 2:
            continue
        direction = g.get("direction", "horizontal")
        axis = "height" if direction == "horizontal" else "width"
        # Only compare group children (skip leaf nodes)
        group_children = [(cid, c) for cid, c in child_bboxes if cid in groups_out]
        if len(group_children) < 2:
            continue
        sizes = [(cid, c[axis]) for cid, c in group_children]
        max_cid, max_s = max(sizes, key=lambda x: x[1])
        min_cid, min_s = min(sizes, key=lambda x: x[1])
        if min_s <= 0 or max_s / min_s < 2.0:
            continue
        max_label = all_items.get(max_cid, {}).get("label", max_cid)
        min_label = all_items.get(min_cid, {}).get("label", min_cid)
        ratio = max_s / min_s
        # Add packing efficiency as supplementary info
        pad = g.get("_padding", {})
        content_w = g["width"] - pad.get("left", 0) - pad.get("right", 0)
        content_h = g["height"] - pad.get("top", 0) - pad.get("bottom", 0)
        content_area = max(content_w, 1) * max(content_h, 1)
        child_area = sum(c["width"] * c["height"] for _, c in child_bboxes)
        eff = round(child_area / content_area * 100)
        warnings.append(f"Group \"{g['label']}\" children {axis} imbalance: \"{max_label}\"={max_s}px vs \"{min_label}\"={min_s}px (ratio {ratio:.1f}:1, packing {eff}%). Consider redistributing children or changing direction. Note: restructuring may affect arrow routing.")

    if warnings:
        output["warnings"] = warnings

    if args.output:
        out_path = Path(args.output)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        write_json(out_path, {"elements": elements})
        print(f"Generated: {out_path}")
        print(f"bbox: x={rb[0]} y={rb[1]} width={rb[2]} height={rb[3]}")
        for w in warnings:
            print(f"⚠️  {w}")
    else:
        print(json.dumps(output, indent=2, ensure_ascii=False))




def cmd_analyze_template(args):
    """Analyze a PPTX template: extract layouts, theme, and placeholder details."""
    from sdpm.analyzer import analyze_template, get_layout_placeholders

    template_path = Path(args.input).resolve()
    if not template_path.exists():
        name = args.input
        if not name.endswith(".pptx"):
            name += ".pptx"
        candidate = Path(__file__).resolve().parent.parent / "templates" / name
        if candidate.exists():
            template_path = candidate
        else:
            print(f"Error: File not found: {args.input}", file=sys.stderr)
            sys.exit(1)

    # --layout: show placeholder details for a specific layout
    if args.layout:
        detail = get_layout_placeholders(template_path, args.layout)
        if not detail:
            print(f"Layout not found: {args.layout}", file=sys.stderr)
            sys.exit(1)
        print(f"layout: {detail['name']}")
        if detail.get("notes"):
            print(f"notes: {detail['notes']}")
        print()
        print("placeholders:")
        for ph in detail["placeholders"]:
            idx = f'"{ph["idx"]}"'
            pos = f'({ph["x"]}, {ph["y"]})'
            size = f'{ph["width"]}x{ph["height"]}'
            fs = ph.get("fontSize")
            fs_str = f'{fs:g} pt ({fs * 2:g} px)' if fs else ""
            desc = ph.get("description", "")
            print(f"  {idx:<5} {pos:<14} {size:<10} {fs_str:<16} {desc}")
        return

    # Full analysis
    result = analyze_template(template_path)

    # Generate color usage and cache preview PNGs if not cached
    if not result["color_usage"]:
        import shutil

        from sdpm.analyzer import cache_color_usage, cache_preview_pngs, extract_color_usage_from_pngs
        preview_dir = Path("/tmp/pptx-preview")
        if preview_dir.exists():
            shutil.rmtree(preview_dir)
        try:
            import contextlib
            import io
            preview_args = type('Args', (), {
                'input': str(template_path), 'pages': None, 'no_grid': True,
                'no_open': True, '_output_dir': str(preview_dir)
            })()
            with contextlib.redirect_stdout(io.StringIO()):
                cmd_preview(preview_args)
        except Exception:
            pass
        usage = extract_color_usage_from_pngs(preview_dir)
        if usage:
            cache_color_usage(template_path, usage)
            result["color_usage"] = usage
        cache_preview_pngs(template_path, preview_dir)

    print(f"template: {template_path.name}")
    sz = result["slide_size"]
    print(f"slideSize: {sz['width']}x{sz['height']}")
    fonts = result.get("fonts", {})
    print("fonts:")
    print(f"  fullwidth: {fonts.get('fullwidth', 'N/A')}")
    print(f"  halfwidth: {fonts.get('halfwidth', 'N/A')}")

    print()
    print("themeColors:")
    for role, color in result["theme_colors"].items():
        print(f"  {role:<15} {color}")

    if result["color_usage"]:
        print()
        print("colorUsage:")
        top5 = result["color_usage"][:5]
        for c in top5:
            print(f"  {c['color']}  {c['percentage']:5.1f}%")
        rest = sum(c["percentage"] for c in result["color_usage"][5:])
        if rest > 0:
            print(f"  other    {rest:5.1f}%")

    print()
    print("layouts:")
    for layout in result["layouts"]:
        print(f'  "layout": "{layout["name"]}"')
        if layout.get("notes"):
            print(f"  {layout['notes']}")
        print()

    ts = result.get("table_styles", {})
    print("tableStyles:")
    if ts.get("styles"):
        for s in ts["styles"]:
            default_mark = " ★default" if s["name"] == ts.get("default") else ""
            print(f'  "{s["name"]}"{default_mark}')
            print(f'    {s["description"]}')
    else:
        print("  (none — theme-based fallback will be used)")
    print()


def cmd_image_size(args):
    """Show image dimensions and calculate size preserving aspect ratio."""
    from PIL import Image

    p = Path(args.input)
    if not p.is_file():
        print(f"Error: Not found: {p}", file=sys.stderr)
        sys.exit(1)

    try:
        img = Image.open(p)
        w, h = img.size
        ratio = w / h
        if args.width:
            calc_h = round(args.width / ratio)
            print(f"{w}x{h} → width={args.width}, height={calc_h}")
        elif args.height:
            calc_w = round(args.height * ratio)
            print(f"{w}x{h} → width={calc_w}, height={args.height}")
    except Exception as e:
        print(f"Error: {e}", file=sys.stderr)
        sys.exit(1)


def cmd_grid(args):
    """Compute CSS Grid layout coordinates."""
    from sdpm.layout.grid import compute_grid
    from sdpm.utils.io import read_json, write_json

    if args.input == "-":
        spec = json.loads(sys.stdin.read())
    else:
        spec = read_json(Path(args.input))
    result = compute_grid(spec)
    if args.output:
        write_json(Path(args.output), result)
    else:
        print(json.dumps(result, indent=2))


def cmd_diff(args):
    """Compare two slide JSONs (or PPTXs) and show manual edit changes."""
    base = load_slides_json_or_pptx(args.baseline)
    edit = load_slides_json_or_pptx(args.edited)

    base_slides = base.get("slides", [])
    edit_slides = edit.get("slides", [])
    _SKIP_KEYS = {"masterIndex", "_comment"}
    has_diff = False

    alignment = align_slides(base_slides, edit_slides)

    for bi, ei in alignment:
        if bi is None:
            es = edit_slides[ei]
            title = es.get("title", "")
            if isinstance(title, dict):
                title = title.get("text", "")
            print(f'\n=== ADDED slide (edited #{ei+1}) "{title[:40]}" ===')
            print(f"  layout: {es.get('layout')}, elements: {len(es.get('elements', []))}")
            has_diff = True
            continue
        if ei is None:
            bs = base_slides[bi]
            title = bs.get("title", "")
            if isinstance(title, dict):
                title = title.get("text", "")
            print(f'\n=== REMOVED slide (baseline #{bi+1}) "{title[:40]}" ===')
            has_diff = True
            continue

        bs, es = base_slides[bi], edit_slides[ei]
        slide_diffs = []

        for key in ("layout", "title", "notes"):
            bv, ev = bs.get(key), es.get(key)
            if bv != ev and (bv or ev):
                slide_diffs.append(_diff_value(key, bv, ev))

        b_elems = [e for e in bs.get("elements", []) if "_comment" not in e]
        e_elems = [e for e in es.get("elements", []) if "_comment" not in e]

        pairs, added = match_elements(b_elems, e_elems)
        elem_diffs = []

        for bj, ej in pairs:
            be = b_elems[bj]
            if ej is None:
                elem_diffs.append(f"  REMOVED [{bj}] {_elem_id(be)}")
                continue
            ee = e_elems[ej]
            all_keys = sorted(set(list(be.keys()) + list(ee.keys())) - _SKIP_KEYS)
            changes = []
            for key in all_keys:
                bv, ev = be.get(key), ee.get(key)
                if bv == ev:
                    continue
                if bv is None:
                    changes.append(f"+{key}={json.dumps(ev, ensure_ascii=False)[:40]}")
                elif ev is None:
                    changes.append(f"-{key}")
                else:
                    if isinstance(bv, (int, float)) and isinstance(ev, (int, float)) and abs(bv - ev) <= 2:
                        continue
                    changes.append(_diff_value(key, bv, ev))
            if changes:
                elem_diffs.append(f"  [{bj}] {_elem_id(be)}:")
                for c in changes:
                    elem_diffs.append(f"    {c}")

        for ej in added:
            ee = e_elems[ej]
            elem_diffs.append(f"  ADDED {_elem_id(ee)}:")
            elem_diffs.append(f"    {json.dumps(ee, ensure_ascii=False)[:300]}")

        moved = bi != ei
        if slide_diffs or elem_diffs or moved:
            title = bs.get("title", es.get("title", ""))
            if isinstance(title, dict):
                title = title.get("text", "")
            moved_str = f" (moved: #{bi+1}→#{ei+1})" if moved else ""
            print(f'\n=== Slide (baseline #{bi+1} ↔ edited #{ei+1}) "{title[:40]}"{moved_str} ===')
            for d in slide_diffs:
                print(f"  {d}")
            for d in elem_diffs:
                print(d)
            has_diff = True

    if not has_diff:
        print("No differences found.")


def main():
    parser = argparse.ArgumentParser(description="PPTX Builder")
    subparsers = parser.add_subparsers(dest="command", required=True)

    p_gen = subparsers.add_parser("generate", help="Generate PPTX from JSON")
    p_gen.add_argument("input", nargs="?", help="Input JSON file (or - for stdin)")
    p_gen.add_argument("-o", "--output", required=True, help="Output PPTX path")
    p_gen.add_argument("--no-autofit", action="store_true", help="Skip PowerPoint autofit refresh")
    p_gen.add_argument("--no-preview", action="store_true", help="Skip automatic preview after generation")
    p_gen.add_argument("--keep-empty-placeholders", action="store_true", help="Keep empty placeholders visible")

    p_prev = subparsers.add_parser("preview", help="Export slides as PNG images")
    p_prev.add_argument("input", help="Input PPTX file")
    p_prev.add_argument("-p", "--pages", help="Pages to export (e.g. 1,3,5)")
    p_prev.add_argument("--no-grid", action="store_true", help="Disable 5% grid overlay")

    p_search = subparsers.add_parser("search-assets", help="Search assets (icons, images, etc.)")
    p_search.add_argument("query", help="Search keywords (space-separated)")
    p_search.add_argument("-n", "--limit", type=int, default=20, help="Max results (default: 20)")
    p_search.add_argument("-s", "--source", help="Filter by source (e.g. aws, material)")
    p_search.add_argument("-t", "--type", help="Filter by type (e.g. service, resource)")
    p_search.add_argument("--theme", choices=["light", "dark"], help="Filter by theme (light/dark)")

    # Backward-compatible alias
    p_icon = subparsers.add_parser("icon-search", help="Search assets (alias for search-assets)")
    p_icon.add_argument("query", help="Search keywords (space-separated)")
    p_icon.add_argument("-n", "--limit", type=int, default=20, help="Max results (default: 20)")
    p_icon.add_argument("-s", "--source", help="Filter by source (e.g. aws, material)")
    p_icon.add_argument("-t", "--type", help="Filter by type (e.g. service, resource)")
    p_icon.add_argument("--theme", choices=["light", "dark"], help="Filter by theme (light/dark)")

    subparsers.add_parser("list-asset-sources", help="List available asset sources")
    subparsers.add_parser("list-templates", help="List available PPTX templates")

    p_ex = subparsers.add_parser("examples", help="List or show design pattern/component examples")
    p_ex.add_argument("names", nargs="*", help="Example names to show (multiple allowed)")
    p_ex.add_argument("--no-browse", action="store_true", help="Don't open browser for styles")

    p_exs = subparsers.add_parser("search-patterns", help="Search patterns by keywords")
    p_exs.add_argument("query", help="Search keywords (space-separated)")
    p_exs.add_argument("-n", "--limit", type=int, default=5, help="Max results")

    p_wf = subparsers.add_parser("workflows", help="List or show workflow documents")
    p_wf.add_argument("names", nargs="*", help="Workflow names to show (multiple allowed)")

    p_gd = subparsers.add_parser("guides", help="List or show guide documents")
    p_gd.add_argument("names", nargs="*", help="Guide names to show (multiple allowed)")

    p_init = subparsers.add_parser("init", help="Initialize output directory with empty presentation JSON")
    p_init.add_argument("name", nargs="?", help="Presentation name (e.g. 'my-proposal')")
    p_init.add_argument("-o", "--output", help="Output directory (overrides default)")
    p_init.add_argument("--style", help="Style name (copies style HTML as art-direction.html)")

    p_layout = subparsers.add_parser("layout", help="Compute layout coordinates from logical structure JSON")

    p_code = subparsers.add_parser("code-block", help="Generate elements JSON for syntax-highlighted code block")
    p_code.add_argument("input", help="Source code file (or - for stdin)")
    p_code.add_argument("-o", "--output", help="Output elements JSON file (default: stdout)")
    p_code.add_argument("--language", "-l", default="text", help="Language for highlighting (default: text)")
    p_code.add_argument("--x", type=int, default=0, help="X position (px)")
    p_code.add_argument("--y", type=int, default=0, help="Y position (px)")
    p_code.add_argument("--width", type=int, default=800, help="Width (px)")
    p_code.add_argument("--height", type=int, default=300, help="Height (px)")
    p_code.add_argument("--font-size", type=int, default=12, help="Font size (pt, default: 12)")
    p_code.add_argument("--font-family", default="Consolas", help="Font family (default: Consolas)")
    p_code.add_argument("--theme", choices=["dark", "light"], default="dark", help="Theme (default: dark)")
    p_code.add_argument("--no-label", action="store_true", help="Hide language label")
    p_layout.add_argument("input", help="Input JSON file (or - for stdin)")
    p_layout.add_argument("-o", "--output", help="Output elements JSON file (default: stdout)")
    p_layout.add_argument("--x", type=int, default=None, help="Target area X offset (px)")
    p_layout.add_argument("--y", type=int, default=None, help="Target area Y offset (px)")
    p_layout.add_argument("--width", type=int, default=None, help="Target area width (px)")
    p_layout.add_argument("--height", type=int, default=None, help="Target area height (px)")
    p_layout.add_argument("--theme", choices=["dark", "light"], default="dark", help="Theme for box text colors (default: dark)")

    p_diff = subparsers.add_parser("diff", help="Compare two slide JSONs/PPTXs and show changes (for manual edit detection)")
    p_diff.add_argument("baseline", help="Baseline slides JSON or PPTX (original)")
    p_diff.add_argument("edited", help="Edited slides JSON or PPTX (manually edited)")

    p_analyze = subparsers.add_parser("analyze-template", help="Analyze PPTX template: extract layouts and theme")
    p_analyze.add_argument("input", help="Template PPTX file path")
    p_analyze.add_argument("--layout", help="Show placeholder details for a specific layout name")

    p_imgsize = subparsers.add_parser("image-size", help="Calculate image size preserving aspect ratio")
    p_imgsize.add_argument("input", help="Image file path")
    p_imgsize_group = p_imgsize.add_mutually_exclusive_group(required=True)
    p_imgsize_group.add_argument("--width", type=int, help="Target width → calculate height")
    p_imgsize_group.add_argument("--height", type=int, help="Target height → calculate width")

    p_grid = subparsers.add_parser("grid", help="Compute CSS Grid layout coordinates")
    p_grid.add_argument("input", help="Input JSON file (or - for stdin)")
    p_grid.add_argument("-o", "--output", help="Output JSON file (default: stdout)")

    args = parser.parse_args()

    if args.command == "generate":
        cmd_generate(args)
    elif args.command == "preview":
        cmd_preview(args)
    elif args.command == "search-assets":
        cmd_search_assets(args)
    elif args.command == "icon-search":
        cmd_icon_search(args)
    elif args.command == "list-asset-sources":
        cmd_list_asset_sources(args)
    elif args.command == "list-templates":
        cmd_list_templates(args)
    elif args.command == "examples":
        cmd_examples(args)
    elif args.command == "search-patterns":
        cmd_search_patterns(args)
    elif args.command == "workflows":
        cmd_workflows(args)
    elif args.command == "guides":
        cmd_guides(args)
    elif args.command == "init":
        cmd_init(args)
    elif args.command == "layout":
        cmd_layout(args)
    elif args.command == "code-block":
        cmd_code_block(args)
    elif args.command == "diff":
        cmd_diff(args)
    elif args.command == "analyze-template":
        cmd_analyze_template(args)
    elif args.command == "image-size":
        cmd_image_size(args)
    elif args.command == "grid":
        cmd_grid(args)

if __name__ == "__main__":
    main()
