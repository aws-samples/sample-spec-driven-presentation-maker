# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Hand-edit sync (Workflow C) tests: deck → generate → hand-edit → diff.

Guards the diff pipeline that syncs manual PowerPoint edits back into the
deck's slide JSON:
- ``load_slides_json_or_pptx`` accepts a deck DIRECTORY as baseline
  (post pptx-import-edit, decks are directories — not a single slides.json)
- ``diff_report`` detects text edits and added elements in the edited PPTX
"""

import json
from pathlib import Path

import pytest

from sdpm.api import generate
from sdpm.diff import diff_report


@pytest.fixture()
def deck(tmp_path: Path) -> Path:
    """Minimal two-slide deck directory."""
    deck = tmp_path / "deck"
    (deck / "slides").mkdir(parents=True)
    (deck / "specs").mkdir()
    (deck / "deck.json").write_text(json.dumps({
        "template": "blank-dark.pptx",
        "fonts": {"fullwidth": "Meiryo", "halfwidth": "Arial"},
        "defaultTextColor": "#FFFFFF",
    }))
    (deck / "specs" / "outline.md").write_text("- [intro] intro\n- [detail] detail\n")
    (deck / "slides" / "intro.json").write_text(json.dumps({
        "layout": "Blank",
        "elements": [
            {"type": "textbox", "text": "Original Title", "x": 100, "y": 100,
             "width": 800, "height": 80, "fontSize": 32},
        ],
    }))
    (deck / "slides" / "detail.json").write_text(json.dumps({
        "layout": "Blank",
        "elements": [
            {"type": "textbox", "text": "Detail body", "x": 100, "y": 300,
             "width": 800, "height": 80, "fontSize": 20},
        ],
    }))
    return deck


def _build(deck: Path, out: Path) -> Path:
    result = generate(deck, output_path=out)
    assert result["slide_count"] == 2
    return out


class TestDiffDeckDirBaseline:
    def test_no_edit_reports_no_diff(self, deck: Path, tmp_path: Path) -> None:
        pptx = _build(deck, tmp_path / "out.pptx")
        result = diff_report(deck, pptx)
        assert result["has_diff"] is False
        assert "No differences found." in result["report"]

    def test_hand_edited_text_is_detected(self, deck: Path, tmp_path: Path) -> None:
        pptx = _build(deck, tmp_path / "out.pptx")

        # Simulate a hand edit in PowerPoint: change the title text
        from pptx import Presentation
        prs = Presentation(str(pptx))
        edited = False
        for shape in prs.slides[0].shapes:
            if shape.has_text_frame and "Original Title" in shape.text_frame.text:
                shape.text_frame.paragraphs[0].runs[0].text = "Hand-Edited Title"
                edited = True
        assert edited, "fixture text not found in generated PPTX"
        edited_pptx = tmp_path / "edited.pptx"
        prs.save(str(edited_pptx))

        result = diff_report(deck, edited_pptx)
        assert result["has_diff"] is True
        assert "Hand-Edited Title" in result["report"]

    def test_added_shape_is_detected(self, deck: Path, tmp_path: Path) -> None:
        pptx = _build(deck, tmp_path / "out.pptx")

        from pptx import Presentation
        from pptx.util import Emu
        prs = Presentation(str(pptx))
        slide = prs.slides[1]
        from pptx.enum.shapes import MSO_SHAPE
        slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(914400), Emu(914400), Emu(1828800), Emu(914400),
        )
        edited_pptx = tmp_path / "edited.pptx"
        prs.save(str(edited_pptx))

        result = diff_report(deck, edited_pptx)
        assert result["has_diff"] is True
        assert "ADDED" in result["report"]


if __name__ == "__main__":
    pytest.main([__file__, "-v"])
