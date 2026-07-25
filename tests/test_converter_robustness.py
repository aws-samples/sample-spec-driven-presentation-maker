# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Converter robustness tests from real-world PPTX corpus findings.

Each test reproduces a bug found by running pptx_to_json over a corpus of
~140 real presentations (Canva/Google Slides exports, business templates):

1. Chart per-point colors used int dict keys → minimal-mode strip crashed
   with AttributeError on every deck containing such charts.
2. Slides with a duplicate placeholder idx (Google Slides exports emit two
   TITLE/idx=0 shapes on one slide) silently dropped the second shape.
3. Multi-paragraph title placeholders lost every paragraph after the first.
"""

import copy
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

from sdpm.converter.pipeline import pptx_to_json
from sdpm.schema.minimal import _strip_internal_keys


def _template() -> Path:
    return Path(__file__).parent.parent / "skill" / "templates" / "blank-dark.pptx"


def _title_slide(prs):
    """Add a slide using a layout that has a title placeholder."""
    for layout in prs.slide_layouts:
        try:
            slide = prs.slides.add_slide(layout)
        except Exception:
            continue
        if slide.shapes.title is not None:
            return slide
        # remove unusable slide? keep simple: continue scanning
    raise AssertionError("no layout with title placeholder in template")


class TestMinimalStripNonStrKeys:
    def test_int_keys_do_not_crash(self):
        slide = {
            "layout": "Blank",
            "elements": [{
                "type": "chart",
                "chartType": "pie",
                "series": [{"name": "s", "values": [1, 2], "pointColors": {0: "#FF0000", 1: "#00FF00"}}],
                "_internal": "dropme",
            }],
        }
        out = _strip_internal_keys(slide)
        el = out["elements"][0]
        assert "_internal" not in el
        # int keys survive the strip untouched (previously: AttributeError)
        assert el["series"][0]["pointColors"] == {0: "#FF0000", 1: "#00FF00"}


class TestChartPointColorKeys:
    def test_extracted_point_colors_use_str_keys(self, tmp_path):
        # Build a pie chart with explicit per-point colors via raw dPt XML
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from lxml import etree

        prs = Presentation(str(_template()))
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        data = CategoryChartData()
        data.categories = ["a", "b"]
        data.add_series("s1", (1.0, 2.0))
        gf = slide.shapes.add_chart(
            XL_CHART_TYPE.PIE, Emu(0), Emu(0), Emu(3000000), Emu(3000000), data
        )
        ser = gf.chart._chartSpace.findall(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/chart}ser")[0]
        ns_c = "http://schemas.openxmlformats.org/drawingml/2006/chart"
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        for i, color in enumerate(["FF0000", "00FF00"]):
            dpt = etree.SubElement(ser, f"{{{ns_c}}}dPt")
            idx = etree.SubElement(dpt, f"{{{ns_c}}}idx")
            idx.set("val", str(i))
            sppr = etree.SubElement(dpt, f"{{{ns_c}}}spPr")
            fill = etree.SubElement(sppr, f"{{{ns_a}}}solidFill")
            srgb = etree.SubElement(fill, f"{{{ns_a}}}srgbClr")
            srgb.set("val", color)
        pptx_path = tmp_path / "chart.pptx"
        prs.save(str(pptx_path))

        # minimal=True used to crash (int keys hit str.startswith in strip)
        result = pptx_to_json(pptx_path, tmp_path / "out", minimal=True)

        charts = [el for s in result["slides"] for el in s.get("elements", []) if el.get("type") == "chart"]
        assert charts, "chart element not extracted"
        pc = charts[0]["series"][0].get("pointColors")
        assert pc, "pointColors not extracted"
        assert all(isinstance(k, str) for k in pc), f"non-str keys: {pc}"


class TestDuplicatePlaceholderIdx:
    def test_second_same_idx_placeholder_is_rescued(self, tmp_path):
        prs = Presentation(str(_template()))
        slide = _title_slide(prs)
        title = slide.shapes.title
        title.text_frame.text = "first title"
        # Clone the title sp element (same placeholder idx) — as emitted by
        # Google Slides exports — and give it different text/position.
        clone = copy.deepcopy(title._element)
        title._element.addnext(clone)
        second = [sh for sh in slide.shapes if sh.shape_id != title.shape_id and sh.is_placeholder][0]
        second.text_frame.text = "second duplicate"
        second.left = Emu(1000000)
        second.top = Emu(2000000)
        pptx_path = tmp_path / "dup.pptx"
        prs.save(str(pptx_path))

        result = pptx_to_json(pptx_path, tmp_path / "out")
        dumped = json.dumps(result["slides"], ensure_ascii=False)
        assert "first title" in dumped
        assert "second duplicate" in dumped, "duplicate-idx placeholder text was dropped"


class TestMultiParagraphTitle:
    def test_title_keeps_all_paragraphs(self, tmp_path):
        prs = Presentation(str(_template()))
        slide = _title_slide(prs)
        tf = slide.shapes.title.text_frame
        tf.text = "line one"
        p2 = tf.add_paragraph()
        p2.text = "line two"
        pptx_path = tmp_path / "multi.pptx"
        prs.save(str(pptx_path))

        result = pptx_to_json(pptx_path, tmp_path / "out")
        dumped = json.dumps(result["slides"], ensure_ascii=False)
        assert "line one" in dumped
        assert "line two" in dumped, "second title paragraph was dropped"


class TestHiddenSlides:
    def test_hidden_flag_roundtrips(self, tmp_path):
        prs = Presentation(str(_template()))
        s1 = prs.slides.add_slide(prs.slide_layouts[0])
        s2 = prs.slides.add_slide(prs.slide_layouts[0])
        s2._element.set("show", "0")  # PowerPoint: Hide Slide
        del s1  # only to silence linters
        pptx_path = tmp_path / "hidden.pptx"
        prs.save(str(pptx_path))

        result = pptx_to_json(pptx_path, tmp_path / "out")
        # template may ship with its own slides — check the last two we added
        hidden_flags = [sl.get("hidden") for sl in result["slides"]]
        assert hidden_flags[-1] is True, f"hidden flag not extracted: {hidden_flags}"
        assert hidden_flags[-2] is not True

        # Builder restores the flag
        from sdpm.builder import PPTXBuilder
        builder = PPTXBuilder(str(_template()), fonts={"fullwidth": "Meiryo", "halfwidth": "Arial"}, default_text_color="#FFFFFF")
        builder.add_slide({"layout": builder_first_layout(builder)})
        builder.add_slide({"layout": builder_first_layout(builder), "hidden": True})
        out = tmp_path / "rebuilt.pptx"
        builder.save(str(out))
        prs2 = Presentation(str(out))
        slides = list(prs2.slides)
        assert slides[0]._element.get("show") != "0"
        assert slides[1]._element.get("show") == "0", "hidden flag not rebuilt"


def builder_first_layout(builder) -> str:
    return next(iter(builder.layouts))


class TestDiffDetectsPlaceholderEdits:
    def test_title_placeholder_edit_is_reported(self, tmp_path):
        from sdpm.diff import diff_report

        base = {"slides": [{"layout": "L", "placeholders": {"0": "original title"}, "elements": []}]}
        edit = {"slides": [{"layout": "L", "placeholders": {"0": "edited title"}, "elements": []}]}
        bp = tmp_path / "base.json"
        ep = tmp_path / "edit.json"
        bp.write_text(json.dumps(base))
        ep.write_text(json.dumps(edit))

        rep = diff_report(bp, ep)
        assert rep["has_diff"], "placeholder edit went undetected"
        assert "placeholder[0]" in rep["report"]


class TestExplicitCropKeepsFrame:
    def test_cropped_image_fills_declared_box(self, tmp_path):
        """With an explicit crop, the frame must NOT shrink to the uncropped
        image's aspect ratio (PowerPoint srcRect semantics: crop fills frame)."""
        from PIL import Image as PILImage

        img = tmp_path / "tall.png"
        PILImage.new("RGB", (400, 800), "red").save(img)  # tall image

        deck = tmp_path / "deck"
        (deck / "slides").mkdir(parents=True)
        (deck / "specs").mkdir()
        (deck / "deck.json").write_text(json.dumps({
            "template": str(_template()),
            "fonts": {"fullwidth": "Meiryo", "halfwidth": "Arial"},
        }))
        (deck / "specs" / "outline.md").write_text("- [s1] test\n")
        (deck / "slides" / "s1.json").write_text(json.dumps({
            "layout": "Blank",
            "elements": [{
                "type": "image", "src": str(img),
                "x": 0, "y": 0, "width": 800, "height": 200,   # wide box
                "crop": {"top": 40.0, "bottom": 35.0},          # crops to wide region
            }],
        }))
        from sdpm.api import generate
        out = tmp_path / "out.pptx"
        generate(deck, output_path=out)

        prs = Presentation(str(out))
        pics = [sh for sh in prs.slides[0].shapes if sh.shape_type == 13]
        assert pics, "picture not built"
        pic = pics[0]
        # Frame keeps the declared 800px-wide box (compare as a fraction of
        # slide width to stay independent of the builder's px scale). The old
        # fit=contain behaviour shrank width to height*aspect = 100px (~5%).
        ratio = pic.width / prs.slide_width
        assert ratio > 0.35, f"frame shrank despite explicit crop: width ratio={ratio:.2f}"


class TestTableThemeStyle:
    def test_table_style_id_resolves_to_style_dict(self, tmp_path):
        """Tables styled via tableStyleId (theme table styles) must emit a
        table-level style dict; otherwise the builder overwrites the look
        with its own default banding (white rows on dark decks)."""
        import zipfile

        from pptx.util import Inches

        prs = Presentation(str(_template()))
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        gf = slide.shapes.add_table(3, 2, Inches(1), Inches(1), Inches(6), Inches(2))
        tbl = gf.table
        for r in range(3):
            for c in range(2):
                tbl.cell(r, c).text = f"r{r}c{c}"
        # style id python-pptx stamped on the table
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        sid = tbl._tbl.find(f"{{{ns_a}}}tblPr/{{{ns_a}}}tableStyleId").text
        raw_path = tmp_path / "raw.pptx"
        prs.save(str(raw_path))

        # Inject a tableStyles.xml defining that style (fill + border)
        table_styles = f'''<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<a:tblStyleLst xmlns:a="{ns_a}" def="{sid}">
 <a:tblStyle styleId="{sid}" styleName="Test Style">
  <a:wholeTbl>
   <a:tcStyle>
    <a:tcBdr><a:insideH><a:ln w="12700"><a:solidFill><a:srgbClr val="AD5CFF"/></a:solidFill></a:ln></a:insideH></a:tcBdr>
    <a:fill><a:solidFill><a:srgbClr val="112233"/></a:solidFill></a:fill>
   </a:tcStyle>
  </a:wholeTbl>
 </a:tblStyle>
</a:tblStyleLst>'''
        pptx_path = tmp_path / "styled_table.pptx"
        with zipfile.ZipFile(raw_path) as zin, zipfile.ZipFile(pptx_path, "w") as zout:
            for item in zin.namelist():
                if item == "ppt/tableStyles.xml":
                    zout.writestr(item, table_styles)
                else:
                    zout.writestr(item, zin.read(item))

        result = pptx_to_json(pptx_path, tmp_path / "out")
        tables = [el for s in result["slides"] for el in s.get("elements", []) if el.get("type") == "table"]
        assert tables, "table element not extracted"
        style = tables[0].get("style")
        assert style is not None, "tableStyleId did not resolve to a style dict"
        assert style["body"]["background"] == "#112233"
        assert style["border"]["color"] == "#AD5CFF"


class TestBackgroundFills:
    def test_gradient_background_becomes_fullslide_rect(self, tmp_path):
        """Slides with a gradient background must not silently lose it —
        the converter emits a full-slide gradient rectangle."""
        from lxml import etree

        prs = Presentation(str(_template()))
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        bg = etree.SubElement(slide._element.find(f"{{{ns_p}}}cSld"), f"{{{ns_p}}}bg")
        bgpr = etree.SubElement(bg, f"{{{ns_p}}}bgPr")
        grad = etree.fromstring(
            f'<a:gradFill xmlns:a="{ns_a}"><a:gsLst>'
            f'<a:gs pos="0"><a:srgbClr val="00FFCC"/></a:gs>'
            f'<a:gs pos="100000"><a:srgbClr val="0066FF"/></a:gs>'
            f'</a:gsLst><a:lin ang="0" scaled="1"/></a:gradFill>')
        bgpr.append(grad)
        etree.SubElement(bgpr, f"{{{ns_a}}}effectLst")
        # move bg before spTree (schema order)
        csld = slide._element.find(f"{{{ns_p}}}cSld")
        csld.remove(bg)
        csld.insert(0, bg)
        pptx_path = tmp_path / "gradbg.pptx"
        prs.save(str(pptx_path))

        result = pptx_to_json(pptx_path, tmp_path / "out")
        last = result["slides"][-1]
        els = last.get("elements", [])
        assert els and els[0].get("type") == "shape" and els[0].get("shape") == "rectangle" \
            and els[0].get("gradient"), f"gradient background not extracted: {els[:1]}"

    def test_image_background_becomes_fullslide_image(self, tmp_path):
        """Slides with a picture-fill background keep it as a cover image."""
        from lxml import etree
        from PIL import Image as PILImage

        img_file = tmp_path / "bg.png"
        PILImage.new("RGB", (32, 18), "blue").save(img_file)

        prs = Presentation(str(_template()))
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        # add picture to obtain an image part + rId, then reference it from bg
        pic = slide.shapes.add_picture(str(img_file), 0, 0)
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        rid = pic._element.find(f".//{{{ns_a}}}blip").get(f"{{{ns_r}}}embed")
        slide.shapes._spTree.remove(pic._element)  # picture itself not needed
        csld = slide._element.find(f"{{{ns_p}}}cSld")
        bg = etree.fromstring(
            f'<p:bg xmlns:p="{ns_p}" xmlns:a="{ns_a}" xmlns:r="{ns_r}"><p:bgPr>'
            f'<a:blipFill><a:blip r:embed="{rid}"/><a:stretch><a:fillRect/></a:stretch></a:blipFill>'
            f'<a:effectLst/></p:bgPr></p:bg>')
        csld.insert(0, bg)
        pptx_path = tmp_path / "imgbg.pptx"
        prs.save(str(pptx_path))

        result = pptx_to_json(pptx_path, tmp_path / "out")
        last = result["slides"][-1]
        els = last.get("elements", [])
        assert els and els[0].get("type") == "image" and "_bg" in els[0].get("src", ""), \
            f"image background not extracted: {els[:1]}"


class TestHideMasterShapes:
    def test_show_master_sp_roundtrips(self, tmp_path):
        """showMasterSp="0" (Hide Background Graphics) must roundtrip;
        losing it lets master decoration cover the slide's own background."""
        prs = Presentation(str(_template()))
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide._element.set("showMasterSp", "0")
        pptx_path = tmp_path / "hidemaster.pptx"
        prs.save(str(pptx_path))

        result = pptx_to_json(pptx_path, tmp_path / "out")
        assert result["slides"][-1].get("hideMasterShapes") is True, \
            "showMasterSp=0 not extracted"

        from sdpm.builder import PPTXBuilder
        builder = PPTXBuilder(str(_template()), fonts={"fullwidth": "Meiryo", "halfwidth": "Arial"}, default_text_color="#FFFFFF")
        builder.add_slide({"layout": builder_first_layout(builder), "hideMasterShapes": True})
        out = tmp_path / "rebuilt.pptx"
        builder.save(str(out))
        prs2 = Presentation(str(out))
        assert list(prs2.slides)[-1]._element.get("showMasterSp") == "0", \
            "showMasterSp=0 not rebuilt"


class TestSrgbColorTransforms:
    def test_lummod_on_srgb_cell_fill(self):
        """<a:srgbClr val="4F81BD"><a:lumMod val="75000"/> must darken the
        color — dropping the transform rendered strong blue instead of the
        original muted tone."""
        from lxml import etree
        from sdpm.converter.color import apply_element_transforms

        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        el = etree.fromstring(
            f'<a:srgbClr xmlns:a="{ns}" val="4F81BD"><a:lumMod val="75000"/></a:srgbClr>')
        out = apply_element_transforms("#4F81BD", el)
        assert out != "#4F81BD"
        # darker than the base
        assert int(out[1:3], 16) < 0x4F


class TestPresetColor:
    def test_prstclr_white_extracted(self, tmp_path):
        from lxml import etree

        prs = Presentation(str(_template()))
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        from pptx.util import Emu
        tb = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(2000000), Emu(500000))
        tb.text_frame.text = "PresetRed"
        run = tb.text_frame.paragraphs[0].runs[0]
        rPr = run._r.get_or_add_rPr()
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        fill = etree.SubElement(rPr, f"{{{ns}}}solidFill")
        etree.SubElement(fill, f"{{{ns}}}prstClr").set("val", "red")
        rPr.insert(0, fill)
        pptx_path = tmp_path / "prst.pptx"
        prs.save(str(pptx_path))

        result = pptx_to_json(pptx_path, tmp_path / "out")
        dumped = json.dumps(result["slides"], ensure_ascii=False)
        assert "#FF0000" in dumped and "PresetRed" in dumped


class TestCellVerticalAlignDefault:
    def test_missing_anchor_becomes_top(self, tmp_path):
        """OOXML default cell anchor is top; the builder's default is middle,
        so the converter must emit vertical-align explicitly."""
        from pptx.util import Inches

        prs = Presentation(str(_template()))
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        gf = slide.shapes.add_table(2, 1, Inches(1), Inches(1), Inches(4), Inches(2))
        gf.table.cell(0, 0).text = "a"
        gf.table.cell(1, 0).text = "b"
        # remove anchor attribute if python-pptx set one
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        for tc in gf.table._tbl.iter(f"{{{ns}}}tcPr"):
            tc.attrib.pop("anchor", None)
        pptx_path = tmp_path / "table_anchor.pptx"
        prs.save(str(pptx_path))

        result = pptx_to_json(pptx_path, tmp_path / "out")
        tables = [el for s in result["slides"] for el in s.get("elements", []) if el.get("type") == "table"]
        cell = tables[0]["rows"][0][0]
        assert isinstance(cell, dict) and cell.get("vertical-align") == "top"


class TestRawShapePassthrough:
    def test_wordart_text_warp_roundtrips(self, tmp_path):
        """prstTxWarp (WordArt arch/circle text) can't be expressed in the
        JSON schema — the shape must roundtrip as rawShape XML."""
        from lxml import etree

        prs = Presentation(str(_template()))
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        from pptx.util import Emu
        tb = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(3000000), Emu(3000000))
        tb.text_frame.text = "ARCHED"
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        bodyPr = tb.text_frame._txBody.find(f"{{{ns}}}bodyPr")
        warp = etree.SubElement(bodyPr, f"{{{ns}}}prstTxWarp")
        warp.set("prst", "textArchUp")
        pptx_path = tmp_path / "wordart.pptx"
        prs.save(str(pptx_path))

        result = pptx_to_json(pptx_path, tmp_path / "out")
        raws = [el for s in result["slides"] for el in s.get("elements", []) if el.get("type") == "rawShape"]
        assert raws and "textArchUp" in raws[0].get("_shapeXml", ""), "WordArt not captured as rawShape"

        # builder injects it back
        from sdpm.builder import PPTXBuilder
        builder = PPTXBuilder(str(_template()), fonts={"fullwidth": "Meiryo", "halfwidth": "Arial"}, default_text_color="#FFFFFF")
        builder.add_slide({"layout": builder_first_layout(builder), "elements": [raws[0]]})
        out = tmp_path / "rebuilt.pptx"
        builder.save(str(out))
        import zipfile
        z = zipfile.ZipFile(str(out))
        slides_xml = "".join(z.read(n).decode() for n in z.namelist() if n.startswith("ppt/slides/slide"))
        assert "textArchUp" in slides_xml, "rawShape not injected on rebuild"


class TestGroupImageReattach:
    def test_group_with_picture_keeps_image(self, tmp_path):
        """Images inside raw-XML groups must be re-attached on rebuild (the
        source rIds dangle in the new package otherwise)."""
        from PIL import Image as PILImage
        from pptx.util import Emu

        img = tmp_path / "photo.png"
        PILImage.new("RGB", (60, 40), "magenta").save(img)

        prs = Presentation(str(_template()))
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        pic = slide.shapes.add_picture(str(img), Emu(0), Emu(0))
        tb = slide.shapes.add_textbox(Emu(0), Emu(500000), Emu(1000000), Emu(300000))
        tb.text_frame.text = "label"
        # freeform member forces _groupXml passthrough; use a real group
        grp = slide.shapes.add_group_shape([pic, tb])
        # inject a fake freeform marker: easier — rotate group to force rawXml path
        grp.rotation = 15.0
        pptx_path = tmp_path / "grouped.pptx"
        prs.save(str(pptx_path))

        result = pptx_to_json(pptx_path, tmp_path / "out")
        groups = [el for s in result["slides"] for el in s.get("elements", []) if el.get("type") == "group"]
        assert groups and groups[0].get("_groupXml")
        assert groups[0].get("_groupImages"), "group-referenced image not saved"
        rel = list(groups[0]["_groupImages"].values())[0]
        assert (tmp_path / "out" / rel).exists()


class TestRotatedConnectors:
    """Connectors carrying xfrm rot= (90/180/270°) lost their orientation.

    The schema has no rotation on line elements, so the converter must bake
    the rotation into the endpoints — and flag 90/270° bent connectors as
    V-H-V via elbowStart so the builder reconstructs the elbow correctly.
    """

    @staticmethod
    def _connector_slide(rot_deg, prst="bentConnector3", flip_v=False):
        import tempfile

        from pptx.enum.shapes import MSO_CONNECTOR

        prs = Presentation(str(_template()))
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.ELBOW, Emu(1270000), Emu(635000), Emu(2540000), Emu(1905000))
        ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        sp_pr = conn._element.spPr
        sp_pr.find(f"{ns}prstGeom").set("prst", prst)
        xfrm = sp_pr.find(f"{ns}xfrm")
        if rot_deg:
            xfrm.set("rot", str(rot_deg * 60000))
        if flip_v:
            xfrm.set("flipV", "1")
        with tempfile.TemporaryDirectory() as td:
            p = Path(td) / "t.pptx"
            prs.save(p)
            result = pptx_to_json(p, Path(td) / "out")
        lines = [e for s in result["slides"] for e in s["elements"] if e.get("type") == "line"]
        assert lines, "connector was not extracted"
        return lines[0]

    def test_180_rotation_swaps_endpoints(self):
        base = self._connector_slide(rot_deg=0)
        rot = self._connector_slide(rot_deg=180)
        # 180° = point reflection through the box center: endpoints swap.
        assert (rot["x1"], rot["y1"]) == (base["x2"], base["y2"])
        assert (rot["x2"], rot["y2"]) == (base["x1"], base["y1"])
        assert "elbowStart" not in rot  # still starts horizontal

    def test_90_rotation_marks_vertical_elbow(self):
        rot = self._connector_slide(rot_deg=90)
        assert rot["elbowStart"] == "vertical"

    def test_270_rotation_marks_vertical_elbow(self):
        rot = self._connector_slide(rot_deg=270, flip_v=True)
        assert rot["elbowStart"] == "vertical"

    def test_unrotated_connector_unchanged(self):
        base = self._connector_slide(rot_deg=0)
        assert "elbowStart" not in base
        assert (base["x1"], base["y1"]) == (200, 100)
        assert (base["x2"], base["y2"]) == (400, 300)


class TestTextboxVerticalAnchor:
    """Textboxes with bodyPr anchor= lost their vertical alignment.

    extract_shape_element extracted verticalAlign but the textbox path did
    not, so center/bottom-anchored labels rendered top-aligned on rebuild
    (builder textbox default is top).
    """

    @staticmethod
    def _textbox_slide(anchor=None):
        import tempfile

        from pptx.util import Emu as E

        prs = Presentation(str(_template()))
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        tb = slide.shapes.add_textbox(E(1270000), E(635000), E(2540000), E(635000))
        tb.text_frame.text = "label"
        if anchor:
            ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
            tb.text_frame._txBody.find(f"{ns}bodyPr").set("anchor", anchor)
        with tempfile.TemporaryDirectory() as td:
            p = Path(td) / "t.pptx"
            prs.save(p)
            result = pptx_to_json(p, Path(td) / "out")
        boxes = [e for s in result["slides"] for e in s["elements"] if e.get("type") == "textbox"]
        assert boxes
        return boxes[0]

    def test_center_anchor_extracted(self):
        assert self._textbox_slide("ctr")["verticalAlign"] == "middle"

    def test_bottom_anchor_extracted(self):
        assert self._textbox_slide("b")["verticalAlign"] == "bottom"

    def test_no_anchor_stays_unset(self):
        assert "verticalAlign" not in self._textbox_slide(None)


class TestSysClrResolution:
    """<a:sysClr> fills/lines resolved to 'none', dropping visible borders."""

    @staticmethod
    def _shape_slide(fill_xml=None, line_xml=None):
        import tempfile

        from lxml import etree
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Emu as E

        prs = Presentation(str(_template()))
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(1270000), E(635000), E(2540000), E(1270000))
        ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
        sp_pr = sp._element.spPr
        # Drop the default style-based fill python-pptx leaves in place
        for tag in ("solidFill", "ln"):
            for el in sp_pr.findall(f"{{{ns}}}{tag}"):
                sp_pr.remove(el)
        style = sp._element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}style")
        if style is not None:
            sp._element.remove(style)
        geom = sp_pr.find(f"{{{ns}}}prstGeom")
        if fill_xml:
            geom.addnext(etree.fromstring(fill_xml))
        if line_xml:
            sp_pr.append(etree.fromstring(line_xml))
        with tempfile.TemporaryDirectory() as td:
            p = Path(td) / "t.pptx"
            prs.save(p)
            result = pptx_to_json(p, Path(td) / "out")
        shapes = [e for s in result["slides"] for e in s["elements"] if e.get("type") == "shape"]
        assert shapes
        return shapes[0]

    A = 'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main"'

    def test_sysclr_line_uses_lastclr(self):
        el = self._shape_slide(line_xml=(
            f'<a:ln {self.A} w="19050"><a:solidFill>'
            '<a:sysClr val="windowText" lastClr="000000"/></a:solidFill></a:ln>'))
        assert el["line"] == "#000000"

    def test_sysclr_fill_applies_lum_transforms_and_alpha(self):
        el = self._shape_slide(fill_xml=(
            f'<a:solidFill {self.A}><a:sysClr val="windowText" lastClr="000000">'
            '<a:lumMod val="20000"/><a:lumOff val="80000"/><a:alpha val="30000"/>'
            '</a:sysClr></a:solidFill>'))
        # 20% lum + 80% offset of black = #CCCCCC light gray
        assert el["fill"] == "#CCCCCC"
        assert el["opacity"] == 0.3

    def test_sysclr_gradient_stops_survive(self):
        """sysClr gradient stops were dropped — a 4-stop green→white
        gradient collapsed to a single green stop (rendered solid)."""
        el = self._shape_slide(fill_xml=(
            f'<a:gradFill {self.A}><a:gsLst>'
            '<a:gs pos="0"><a:srgbClr val="00B050"/></a:gs>'
            '<a:gs pos="37000"><a:sysClr val="window" lastClr="FFFFFF"/></a:gs>'
            '<a:gs pos="100000"><a:sysClr val="window" lastClr="FFFFFF"/></a:gs>'
            '</a:gsLst><a:lin ang="0" scaled="1"/></a:gradFill>'))
        stops = el["gradient"]["stops"]
        assert [(s["position"], s["color"]) for s in stops] == [
            (0.0, "#00B050"), (0.37, "#FFFFFF"), (1.0, "#FFFFFF")]


class TestSrgbTintSemantics:
    """ECMA-376 tint: val=100000 must leave the color unchanged.

    The old formula was inverted (100% tint → white), washing out
    gradient fills that Office writes with tint/shade/satMod stops.
    """

    def test_full_tint_is_identity(self):
        from sdpm.converter.xml_helpers import _apply_srgb_transforms
        assert _apply_srgb_transforms("#B22600", {"tint": "100000"}) == "#B22600"

    def test_partial_tint_mixes_toward_white(self):
        from sdpm.converter.xml_helpers import _apply_srgb_transforms
        out = _apply_srgb_transforms("#B22600", {"tint": "50000"})
        r, g, b = int(out[1:3], 16), int(out[3:5], 16), int(out[5:7], 16)
        assert r > 0xB2 and g > 0x26 and b > 0x00  # lighter than input
        assert out != "#FFFFFF"

    def test_satmod_boosts_saturation(self):
        from sdpm.converter.xml_helpers import _apply_srgb_transforms
        out = _apply_srgb_transforms("#996666", {"satMod": "300000"})
        r, g, b = int(out[1:3], 16), int(out[3:5], 16), int(out[5:7], 16)
        assert r > g and r > b  # pushed toward pure red
        assert (r - g) > (0x99 - 0x66)  # more separation than input


class TestImportTextFidelity:
    """Header/label text drifted on rebuild (found on a real deck).

    - align="justify" fell back to center in the shape builder
    - CJK↔Latin auto-spacing mutated imported text (autoSpacing deck flag)
    - empty paragraphs lost their endParaRPr size, shifting anchored text
    """

    def test_parse_styled_text_auto_spacing_off(self):
        from sdpm.utils.text import parse_styled_text
        text = "{{font=Meiryo UI:自動化基盤は}}{{bold,font=Meiryo UI:CLAP/TMT/DNA}}{{font=Meiryo UI:により構成される}}"
        on = parse_styled_text(text)
        off = parse_styled_text(text, auto_spacing=False)
        assert on[0]["text"].endswith(" ")  # spacing applied by default
        assert off[0]["text"] == "自動化基盤は"  # verbatim when disabled
        assert off[1]["text"] == "CLAP/TMT/DNA"

    def test_converted_deck_disables_auto_spacing(self):
        import tempfile
        prs = Presentation(str(_template()))
        prs.slides.add_slide(prs.slide_layouts[-1])
        with tempfile.TemporaryDirectory() as td:
            p = Path(td) / "t.pptx"
            prs.save(p)
            pptx_to_json(p, Path(td) / "out")
            deck = json.loads((Path(td) / "out" / "deck.json").read_text())
        assert deck["autoSpacing"] is False

    def test_sized_empty_paragraphs_roundtrip(self):
        import tempfile

        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Emu as E, Pt

        prs = Presentation(str(_template()))
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(1270000), E(635000), E(2540000), E(1905000))
        tf = sp.text_frame
        tf.paragraphs[0].text = "Title"
        tf.paragraphs[0].runs[0].font.size = Pt(12)
        for _ in range(2):
            para = tf.add_paragraph()  # empty spacer, larger than the text
            from lxml import etree
            from pptx.oxml.ns import qn
            end = etree.SubElement(para._p, qn('a:endParaRPr'))
            end.set('sz', '1600')
        with tempfile.TemporaryDirectory() as td:
            p = Path(td) / "t.pptx"
            prs.save(p)
            result = pptx_to_json(p, Path(td) / "out")
        shapes = [e for s in result["slides"] for e in s["elements"]
                  if e.get("type") == "shape" and "paragraphs" in e]
        assert shapes, "sized empty paragraphs should force paragraphs mode"
        paras = shapes[0]["paragraphs"]
        assert [pa.get("endFontSize") for pa in paras[1:]] == [16.0, 16.0]


class TestTextSizingAndWrap:
    """Two more label-fidelity bugs from a real deck.

    - autoWidth (wrap=none) was ignored when the textbox also carried
      _noAutofit, so no-wrap labels re-wrapped mid-word on rebuild
    - runs with no explicit sz inherit the presentation default (spec
      fallback 18pt), but rebuilt text got the builder default (14pt)
    """

    def test_autowidth_survives_noautofit(self):
        import tempfile
        prs = Presentation(str(_template()))
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        from pptx.util import Emu as E
        tb = slide.shapes.add_textbox(E(1270000), E(635000), E(1905000), E(444500))
        tb.text_frame.text = "シナリオ開発プロセス"
        ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        body = tb.text_frame._txBody.find(f"{ns}bodyPr")
        body.set("wrap", "none")
        from lxml import etree
        etree.SubElement(body, f"{ns}noAutofit")
        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "t.pptx"
            prs.save(src)
            result = pptx_to_json(src, Path(td) / "out")
            el = next(e for s in result["slides"] for e in s["elements"] if e.get("type") == "textbox")
            assert el.get("autoWidth") is True
            el["_noAutofit"] = True  # coexists with autoWidth on real decks
            # Rebuild and confirm wrap=none survives
            from sdpm.builder import PPTXBuilder
            b = PPTXBuilder(_template(), fonts={"fullwidth": "Meiryo", "halfwidth": "Arial"},
                            default_text_color="#000000")
            b.add_slide({"layout": "Blank", "elements": [el]})
            out = Path(td) / "rebuilt.pptx"
            b.save(out)
            reprs = Presentation(str(out))
            tb2 = next(sh for sh in reprs.slides[0].shapes if sh.has_text_frame and "シナリオ" in sh.text_frame.text)
            body2 = tb2.text_frame._txBody.find(f"{ns}bodyPr")
            assert body2.get("wrap") == "none"

    def test_unsized_runs_get_inherited_default(self):
        import tempfile

        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Emu as E
        prs = Presentation(str(_template()))
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(1270000), E(635000), E(5080000), E(952500))
        sp.text_frame.text = "運用者がワークフローエディタを用いて作成"
        # no explicit run size — inherits presentation default (18pt fallback)
        with tempfile.TemporaryDirectory() as td:
            src = Path(td) / "t.pptx"
            prs.save(src)
            result = pptx_to_json(src, Path(td) / "out")
        el = next(e for s in result["slides"] for e in s["elements"]
                  if e.get("type") == "shape" and "運用者" in str(e.get("text", "")))
        assert el.get("fontSize") == 18


class TestTableCellFidelity:
    """Table-cell fidelity bugs from a real deck's spec table.

    - per-cell lstStyle text color (white header text) fell back to tx1 black
    - sysClr cell borders lost their color (white gridlines vanished)
    - buChar bullets inside cells were dropped (builder has no cell lists)
    """

    @staticmethod
    def _convert(mutate):
        import tempfile

        from pptx.util import Emu as E

        prs = Presentation(str(_template()))
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        gfx = slide.shapes.add_table(2, 2, E(1270000), E(635000), E(5080000), E(1905000))
        mutate(gfx.table)
        with tempfile.TemporaryDirectory() as td:
            p = Path(td) / "t.pptx"
            prs.save(p)
            result = pptx_to_json(p, Path(td) / "out")
        for s in result["slides"]:
            for e in s["elements"]:
                if e.get("type") == "table":
                    return e
        raise AssertionError("no table extracted")

    NS = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

    def test_cell_lststyle_color_extracted(self):
        from lxml import etree

        def mutate(tbl):
            cell = tbl.rows[0].cells[0]
            cell.text = "概要"
            txBody = cell._tc.find(f"{self.NS}txBody")
            lst = txBody.find(f"{self.NS}lstStyle")
            lvl = etree.SubElement(lst, f"{self.NS}lvl1pPr")
            d = etree.SubElement(lvl, f"{self.NS}defRPr")
            sf = etree.SubElement(d, f"{self.NS}solidFill")
            etree.SubElement(sf, f"{self.NS}srgbClr").set("val", "FFEE00")
        t = self._convert(mutate)
        cell = t["headers"][0] if t.get("headers") else t["rows"][0][0]
        assert isinstance(cell, dict) and cell.get("color") == "#FFEE00"
        assert "#000000" not in str(cell.get("text"))

    def test_sysclr_border_color_extracted(self):
        from lxml import etree

        def mutate(tbl):
            cell = tbl.rows[0].cells[0]
            cell.text = "x"
            tcPr = cell._tc.get_or_add_tcPr()
            ln = etree.SubElement(tcPr, f"{self.NS}lnB")
            ln.set("w", "12700")
            sf = etree.SubElement(ln, f"{self.NS}solidFill")
            sys_el = etree.SubElement(sf, f"{self.NS}sysClr")
            sys_el.set("val", "window")
            sys_el.set("lastClr", "FFFFFF")
        t = self._convert(mutate)
        cell = t["headers"][0] if t.get("headers") else t["rows"][0][0]
        assert cell["borders"]["bottom"]["color"] == "#FFFFFF"

    def test_cell_bullets_roundtrip(self):
        import tempfile

        from lxml import etree

        def mutate(tbl):
            cell = tbl.rows[0].cells[0]
            tf = cell.text_frame
            tf.text = "各システムに対するドメイン知識"
            p2 = tf.add_paragraph()
            p2.text = "CLAPの仕様理解"
            for para in tf.paragraphs:
                pPr = para._element.get_or_add_pPr()
                pPr.set("marL", "93663")
                pPr.set("indent", "-93663")
                bu = etree.SubElement(pPr, f"{self.NS}buChar")
                bu.set("char", "•")
        t = self._convert(mutate)
        cell = t["headers"][0] if t.get("headers") else t["rows"][0][0]
        assert isinstance(cell, dict)
        paras = cell["paragraphs"]
        assert [p.get("bullet") for p in paras] == ["•", "•"]
        assert [p.get("marL") for p in paras] == [93663, 93663]
        # Rebuild: builder must emit real buChar bullets on each a:p
        from sdpm.builder import PPTXBuilder
        b = PPTXBuilder(_template(), fonts={"fullwidth": "Meiryo", "halfwidth": "Arial"},
                        default_text_color="#000000", auto_spacing=False)
        b.add_slide({"layout": "Blank", "elements": [t]})
        with tempfile.TemporaryDirectory() as td:
            out = Path(td) / "rebuilt.pptx"
            b.save(out)
            prs = Presentation(str(out))
            gfx = next(sh for sh in prs.slides[0].shapes if sh.has_table)
        tc = gfx.table.rows[0].cells[0]._tc
        bu_chars = tc.findall(f".//{self.NS}buChar")
        assert [b_.get("char") for b_ in bu_chars] == ["•", "•"]
        pprs = tc.findall(f".//{self.NS}pPr")
        assert [p_.get("marL") for p_ in pprs] == ["93663", "93663"]
        texts = [p.text for p in gfx.table.rows[0].cells[0].text_frame.paragraphs]
        assert texts == ["各システムに対するドメイン知識", "CLAPの仕様理解"]


class TestNoEffectsSuppression:
    """Rebuilt shapes gained a theme shadow the source never had.

    python-pptx's add_shape writes a default <p:style> whose effectRef
    references the theme effect style (a shadow in both built-in
    templates). Shapes converted from decks whose spPr carries no
    effects now emit _noEffects, and the builder drops the style and
    writes an empty effectLst.
    """

    def test_plain_shape_marks_no_effects(self):
        import tempfile

        from pptx.enum.shapes import MSO_SHAPE
        from pptx.util import Emu as E
        prs = Presentation(str(_template()))
        slide = prs.slides.add_slide(prs.slide_layouts[-1])
        sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(1270000), E(635000), E(2540000), E(635000))
        # python-pptx's own style stays (effectRef idx=2 in the template
        # theme), but spPr has no effectLst → still "no effects" per source?
        # No: effectRef>0 means themed effects apply. Clear the style to
        # simulate the common hand-drawn shape (no style, no effects).
        style = sp._element.find("{http://schemas.openxmlformats.org/presentationml/2006/main}style")
        if style is not None:
            sp._element.remove(style)
        with tempfile.TemporaryDirectory() as td:
            p = Path(td) / "t.pptx"
            prs.save(p)
            result = pptx_to_json(p, Path(td) / "out")
        el = next(e for s in result["slides"] for e in s["elements"] if e.get("type") == "shape")
        assert el.get("_noEffects") is True

    def test_builder_suppresses_theme_shadow(self):
        import tempfile

        from sdpm.builder import PPTXBuilder
        b = PPTXBuilder(_template(), fonts={"fullwidth": "Meiryo", "halfwidth": "Arial"},
                        default_text_color="#000000")
        b.add_slide({"layout": "Blank", "elements": [
            {"type": "shape", "shape": "rectangle", "x": 100, "y": 100,
             "width": 300, "height": 100, "fill": "#F2F2F2", "_noEffects": True},
            {"type": "shape", "shape": "rectangle", "x": 100, "y": 300,
             "width": 300, "height": 100, "fill": "#F2F2F2"},
        ]})
        with tempfile.TemporaryDirectory() as td:
            out = Path(td) / "t.pptx"
            b.save(out)
            prs = Presentation(str(out))
        ns_p = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
        ns_a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
        shapes = [sh for sh in prs.slides[0].shapes if sh.shape_type is not None]
        marked, unmarked = shapes[0], shapes[1]
        assert marked._element.find(f"{ns_p}style") is None
        assert marked._element.spPr.find(f"{ns_a}effectLst") is not None
        # default behavior unchanged for AI-authored decks
        assert unmarked._element.find(f"{ns_p}style") is not None


class TestOleSanitization:
    """Injected raw group XML carried embedded OLE objects whose r:id
    pointed at parts that don't exist in the rebuilt package — PowerPoint
    reported the file as damaged and stripped the slide. OLE frames are
    swapped for their mc:Fallback picture."""

    def test_ole_graphicframe_replaced_by_fallback_pic(self):
        from lxml import etree

        from sdpm.builder import PPTXBuilder
        ns_p = "http://schemas.openxmlformats.org/presentationml/2006/main"
        ns_a = "http://schemas.openxmlformats.org/drawingml/2006/main"
        ns_r = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
        xml = f'''<p:grpSp xmlns:p="{ns_p}" xmlns:a="{ns_a}" xmlns:r="{ns_r}"
            xmlns:mc="http://schemas.openxmlformats.org/markup-compatibility/2006">
          <p:graphicFrame>
            <p:xfrm><a:off x="100" y="200"/><a:ext cx="300" cy="400"/></p:xfrm>
            <a:graphic><a:graphicData uri="{ns_p.replace('presentationml/2006/main','presentationml/2006/ole')}">
              <mc:AlternateContent>
                <mc:Choice Requires="v"><p:oleObj r:id="rId99" imgW="1" imgH="1"/></mc:Choice>
                <mc:Fallback><p:oleObj>
                  <p:pic>
                    <p:blipFill><a:blip r:embed="rId98"/></p:blipFill>
                    <p:spPr/>
                  </p:pic>
                </p:oleObj></mc:Fallback>
              </mc:AlternateContent>
            </a:graphicData></a:graphic>
          </p:graphicFrame>
        </p:grpSp>'''
        el = etree.fromstring(xml)
        b = PPTXBuilder(_template(), fonts={"fullwidth": "Meiryo", "halfwidth": "Arial"},
                        default_text_color="#000000")
        b._sanitize_injected_xml(el)
        assert el.find(f".//{{{ns_p}}}oleObj") is None
        assert el.find(f".//{{{ns_p}}}graphicFrame") is None
        pic = el.find(f".//{{{ns_p}}}pic")
        assert pic is not None
        off = pic.find(f".//{{{ns_a}}}xfrm/{{{ns_a}}}off")
        assert off is not None and off.get("x") == "100"
