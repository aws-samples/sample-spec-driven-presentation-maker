# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Tests for mcp-local compose_slide — the CC Phase 2 per-slide isolation tool.

compose_slide builds/renders a single assigned slug in a throwaway one-slide deck
directory (method A'), without touching the shared output.pptx / preview/ or the
.save.lock, and without generating compose/defs_. These tests pin that isolation
contract plus the two method-A regressions (relative images, font token check).
"""

import json
import shutil
import sys
from pathlib import Path

import anyio
import pytest

_ROOT = Path(__file__).resolve().parent.parent
# sandbox_tools / sandbox are unique to mcp-local (no collision with mcp-server).
sys.path.insert(0, str(_ROOT / "mcp-local"))

import sandbox_tools  # noqa: E402

_TEMPLATE = _ROOT / "skill" / "templates" / "blank-light.pptx"

# Rendering (build/preview/measure) needs LibreOffice + poppler. Skip those
# assertions when unavailable so the suite still runs in minimal CI.
_HAS_SOFFICE = shutil.which("soffice") is not None or Path(
    "/Applications/LibreOffice.app/Contents/MacOS/soffice"
).exists()
_HAS_PDFTOPPM = shutil.which("pdftoppm") is not None
_CAN_RENDER = _HAS_SOFFICE and _HAS_PDFTOPPM


def _make_deck(tmp_path: Path, *, with_image: bool = False, fs_tokens: str | None = None) -> Path:
    """Build a minimal deck directory with deck.json + specs + outline."""
    deck = tmp_path / "deck"
    (deck / "slides").mkdir(parents=True)
    (deck / "specs").mkdir(parents=True)
    deck_meta = {
        "template": str(_TEMPLATE.resolve()),
        "fonts": {"fullwidth": "Arial", "halfwidth": "Arial"},
        "defaultTextColor": "#333333",
    }
    (deck / "deck.json").write_text(json.dumps(deck_meta), encoding="utf-8")
    (deck / "specs" / "outline.md").write_text("- [title] Hello\n- [second] World\n", encoding="utf-8")
    if fs_tokens is not None:
        (deck / "specs" / "art-direction.html").write_text(
            f"<html><head><style>:root{{{fs_tokens}}}</style></head><body></body></html>",
            encoding="utf-8",
        )
    if with_image:
        from PIL import Image

        (deck / "images").mkdir()
        Image.new("RGB", (64, 48), (10, 120, 200)).save(deck / "images" / "logo.png")
    return deck


def _compose(deck: Path, slug: str, code: str, measure: bool = False) -> dict:
    # compose_slide is async (offloads blocking work to a thread); drive it
    # from the sync test via anyio.run.
    out = anyio.run(
        lambda: sandbox_tools.compose_slide(
            purpose="test", code=code, deck_id=str(deck), slug=slug, measure=measure
        )
    )
    return json.loads(out)


# --------------------------------------------------------------------------- #
# Input validation / guardrail
# --------------------------------------------------------------------------- #


def test_rejects_missing_deck(tmp_path):
    """Non-directory deck_id returns an error, not a crash."""
    result = json.loads(
        anyio.run(
            lambda: sandbox_tools.compose_slide(
                purpose="x", code="pass", deck_id=str(tmp_path / "nope"), slug="title"
            )
        )
    )
    assert "Error" in result["output"]


def test_rejects_import_code(tmp_path):
    """Sandbox code with import is rejected (same policy as run_python)."""
    deck = _make_deck(tmp_path)
    result = _compose(deck, "title", "import os")
    assert "rejected" in result["output"].lower()


def test_reports_when_slug_not_written(tmp_path):
    """If user code doesn't write slides/{slug}.json, we surface a clear error."""
    deck = _make_deck(tmp_path)
    result = _compose(deck, "title", 'print("did nothing")')
    assert "compose_error" in result
    assert "title" in result["compose_error"]


def test_docstring_has_guardrail():
    """R2: docstring must warn against casual use and point to run_python."""
    doc = sandbox_tools.compose_slide.__doc__
    assert "INTERNAL" in doc
    assert "run_python" in doc


# --------------------------------------------------------------------------- #
# Isolation contract (no soffice needed — checks side effects on the deck dir)
# --------------------------------------------------------------------------- #

_VALID_SLIDE = (
    'write_json("slides/{slug}.json", {{"layout":"Title Only","title":"Hi",'
    '"elements":[{{"type":"textbox","x":100,"y":300,"width":800,"height":200,'
    '"fontSize":40,"text":"Hi"}}]}})'
)


def test_writes_slide_json(tmp_path):
    """R3: the single writer path writes slides/{slug}.json."""
    deck = _make_deck(tmp_path)
    _compose(deck, "title", _VALID_SLIDE.format(slug="title"))
    assert (deck / "slides" / "title.json").exists()


def test_does_not_create_shared_output_pptx(tmp_path):
    """R3: compose_slide must not write the deck-wide output.pptx."""
    deck = _make_deck(tmp_path)
    _compose(deck, "title", _VALID_SLIDE.format(slug="title"))
    assert not (deck / "output.pptx").exists()


def test_does_not_take_save_lock(tmp_path):
    """R7: compose_slide must not create/hold the .save.lock."""
    deck = _make_deck(tmp_path)
    _compose(deck, "title", _VALID_SLIDE.format(slug="title"))
    assert not (deck / ".save.lock").exists()


def test_does_not_generate_compose_defs(tmp_path):
    """R4: no compose/defs_ generation (CC has no consumer)."""
    deck = _make_deck(tmp_path)
    _compose(deck, "title", _VALID_SLIDE.format(slug="title"))
    assert not (deck / "compose").exists()


def test_does_not_wipe_sibling_preview(tmp_path):
    """R3: a sibling composer's preview PNG must survive (no rmtree of preview/)."""
    deck = _make_deck(tmp_path)
    (deck / "preview").mkdir()
    (deck / "preview" / "second.png").write_bytes(b"SIBLING")
    _compose(deck, "title", _VALID_SLIDE.format(slug="title"))
    assert (deck / "preview" / "second.png").read_bytes() == b"SIBLING"


def test_cleans_up_iso_workdir(tmp_path):
    """The throwaway iso directory is removed after the call."""
    deck = _make_deck(tmp_path)
    _compose(deck, "title", _VALID_SLIDE.format(slug="title"))
    work = deck / "_work"
    leftover = list(work.glob("compose-*")) if work.exists() else []
    assert not leftover


# --------------------------------------------------------------------------- #
# Rendering + method-A regressions (need LibreOffice/poppler)
# --------------------------------------------------------------------------- #


@pytest.mark.skipif(not _CAN_RENDER, reason="LibreOffice/poppler required")
def test_returns_slug_unique_preview(tmp_path):
    """R5: preview_files is returned and named preview/{slug}.png (not page01-*)."""
    deck = _make_deck(tmp_path)
    result = _compose(deck, "title", _VALID_SLIDE.format(slug="title"))
    assert result.get("preview_files"), result
    assert Path(result["preview_files"][0]).name == "title.png"
    assert (deck / "preview" / "title.png").exists()


@pytest.mark.skipif(not _CAN_RENDER, reason="LibreOffice required")
def test_font_token_check_fires(tmp_path):
    """Critical-2 regression: a fontSize outside --fs-* tokens is flagged.

    Method A (temp JSON file input) would resolve base_dir to the temp dir,
    fail to find art-direction.html, and silently skip this check.
    """
    deck = _make_deck(tmp_path, fs_tokens="--fs-title:40pt;--fs-body:20pt;")
    code = (
        'write_json("slides/title.json", {"layout":"Title Only","title":"Hi",'
        '"elements":[{"type":"textbox","x":100,"y":300,"width":800,"height":100,'
        '"fontSize":33,"text":"bad size"}]})'
    )
    result = _compose(deck, "title", code)
    warnings = " ".join(result.get("warnings", []))
    assert "token discipline" in warnings, result
    assert "33pt" in warnings, result


@pytest.mark.skipif(not _CAN_RENDER, reason="LibreOffice required")
def test_relative_image_embedded(tmp_path):
    """Critical-1 regression: a relative images/ reference is embedded in the PPTX.

    Method A would set base_dir to the temp dir, so images/logo.png would not be
    found and the picture would be silently dropped. We rebuild the one-slide deck
    the same way compose_slide does (directory input) and assert the picture lands.
    """
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    from sdpm.api import generate

    deck = _make_deck(tmp_path, with_image=True, fs_tokens="--fs-title:40pt;")
    code = (
        'write_json("slides/title.json", {"layout":"Title Only","title":"Pic",'
        '"elements":[{"type":"image","x":100,"y":300,"width":200,"height":150,'
        '"src":"images/logo.png"}]})'
    )
    # Run compose_slide to write + sanitize the slide (and prove it doesn't error).
    result = _compose(deck, "title", code)
    assert "pptx_error" not in result, result

    # Reproduce the iso build to inspect the embedded media.
    iso = tmp_path / "iso"
    (iso / "specs").mkdir(parents=True)
    (iso / "slides").mkdir(parents=True)
    shutil.copy2(deck / "deck.json", iso / "deck.json")
    (iso / "specs" / "outline.md").write_text("- [title] Pic\n", encoding="utf-8")
    shutil.copy2(deck / "specs" / "art-direction.html", iso / "specs" / "art-direction.html")
    shutil.copy2(deck / "slides" / "title.json", iso / "slides" / "title.json")
    import os

    os.symlink(deck / "images", iso / "images")
    pptx = iso / "one.pptx"
    generate(json_path=str(iso), output_path=str(pptx))

    prs = Presentation(str(pptx))
    pictures = sum(
        1 for s in prs.slides for sh in s.shapes if sh.shape_type == MSO_SHAPE_TYPE.PICTURE
    )
    assert pictures >= 1, "relative image was dropped (method A regression)"
